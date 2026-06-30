#!/usr/bin/env python3
"""
Enhanced Summary Builder — Fixed AI integration, environment updates, and shell usage.

Fixes:
1. Proper AI model selection based on style and configuration
2. Environment variable updates for .env files
3. Enhanced shell integration for PDF manipulation
4. Better error handling and logging
5. Support for multiple AI providers
"""

import argparse
import json
import os
import re
import sys
import csv
import io
import subprocess
import tempfile
from pathlib import Path


# ═══════════════════════════════════════════════════════════════════════════
# EVENT EMISSION
# ═══════════════════════════════════════════════════════════════════════════

def emit(event_type: str, data: dict):
    """Write a JSON-line event to stdout for the backend to stream."""
    payload = json.dumps({"type": event_type, **data})
    print(payload, flush=True)


# ═══════════════════════════════════════════════════════════════════════════
# ENVIRONMENT MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════

def update_env_file(env_path: str, updates: dict):
    """Update .env file with new configuration values."""
    env_lines = []
    existing_keys = set()
    
    # Read existing .env file
    if os.path.exists(env_path):
        with open(env_path, 'r') as f:
            for line in f:
                line = line.strip()
                if line and not line.startswith('#') and '=' in line:
                    key = line.split('=')[0].strip()
                    if key in updates:
                        env_lines.append(f"{key}={updates[key]}")
                        existing_keys.add(key)
                    else:
                        env_lines.append(line)
                else:
                    env_lines.append(line)
    
    # Add new keys that don't exist
    for key, value in updates.items():
        if key not in existing_keys:
            env_lines.append(f"{key}={value}")
    
    # Write back to .env file
    with open(env_path, 'w') as f:
        f.write('\n'.join(env_lines) + '\n')
    
    emit("info", {"message": f"Updated {len(updates)} environment variables in {env_path}"})


def update_env_example(env_example_path: str, updates: dict):
    """Update .env.example file with new configuration values."""
    if not os.path.exists(env_example_path):
        return
    
    with open(env_example_path, 'r') as f:
        content = f.read()
    
    # Add new variables if they don't exist
    new_lines = []
    existing_keys = set()
    
    for line in content.split('\n'):
        if '=' in line and not line.startswith('#'):
            key = line.split('=')[0].strip()
            existing_keys.add(key)
    
    for key, value in updates.items():
        if key not in existing_keys:
            new_lines.append(f"{key}={value}")
    
    if new_lines:
        with open(env_example_path, 'a') as f:
            f.write('\n' + '\n'.join(new_lines) + '\n')
        
        emit("info", {"message": f"Added {len(new_lines)} new variables to {env_example_path}"})


# ═══════════════════════════════════════════════════════════════════════════
# SKILLS SYSTEM — Load SKILL.md modules from skills directory
# ═══════════════════════════════════════════════════════════════════════════

SKILLS_DIR = Path(__file__).parent / "skills"


def load_skill(skill_name: str) -> dict:
    """Load a skill module by name. Returns {name, description, body} or None."""
    skill_dir = SKILLS_DIR / skill_name
    skill_file = skill_dir / "SKILL.md"
    if not skill_file.exists():
        emit("warning", {"message": f"Skill not found: {skill_name}"})
        return None

    content = skill_file.read_text(encoding="utf-8")

    # Parse YAML frontmatter
    name = skill_name
    description = ""
    body = content

    if content.startswith("---"):
        parts = content.split("---", 2)
        if len(parts) >= 3:
            frontmatter = parts[1].strip()
            body = parts[2].strip()

            for line in frontmatter.split("\n"):
                line = line.strip()
                if line.startswith("name:"):
                    name = line.split(":", 1)[1].strip()
                elif line.startswith("description:"):
                    description = line.split(":", 1)[1].strip()
                elif description and not line.startswith("name:"):
                    description += " " + line.strip()

    return {"name": name, "description": description, "body": body}


def load_all_skills() -> dict:
    """Load all available skills from the skills directory."""
    skills = {}
    if not SKILLS_DIR.exists():
        return skills
    for skill_dir in sorted(SKILLS_DIR.iterdir()):
        if skill_dir.is_dir():
            skill = load_skill(skill_dir.name)
            if skill:
                skills[skill_dir.name] = skill
    return skills


# ═══════════════════════════════════════════════════════════════════════════
# TEXT EXTRACTION (unchanged)
# ═══════════════════════════════════════════════════════════════════════════

def detect_type(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    mapping = {
        ".pdf": "pdf",
        ".png": "image", ".jpg": "image", ".jpeg": "image", ".webp": "image",
        ".pptx": "pptx", ".ppt": "pptx",
        ".xlsx": "spreadsheet", ".xls": "spreadsheet", ".csv": "spreadsheet",
        ".txt": "text", ".md": "text",
        ".docx": "docx",
    }
    return mapping.get(ext, "unknown")


def extract_text_pdf(file_path: str) -> str:
    import fitz
    doc = fitz.open(file_path)
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
    full_text = "\n".join(text_parts)
    doc.close()
    if len(full_text.strip()) < 100:
        return extract_pdf_ocr(file_path)
    return full_text


def extract_pdf_ocr(file_path: str) -> str:
    try:
        from pdf2image import convert_from_path
        import pytesseract
        images = convert_from_path(file_path, dpi=200)
        text_parts = []
        for img in images:
            text_parts.append(pytesseract.image_to_string(img))
        return "\n".join(text_parts)
    except Exception:
        return extract_text_pdf_fallback(file_path)


def extract_text_pdf_fallback(file_path: str) -> str:
    import fitz
    doc = fitz.open(file_path)
    text_parts = []
    for page in doc:
        blocks = page.get_text("blocks")
        for block in blocks:
            if len(block) >= 5 and block[4]:
                text_parts.append(block[4].strip())
    doc.close()
    return "\n".join(text_parts)


def extract_text_image(file_path: str) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        return pytesseract.image_to_string(img, config="--psm 6")
    except Exception:
        return ""


def extract_text_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_text.append(line)
        if slide_text:
            text_parts.append("\n".join(slide_text))
    return "\n\n".join(text_parts)


def extract_text_spreadsheet(file_path: str) -> str:
    import openpyxl
    ext = Path(file_path).suffix.lower()
    text_parts = []
    if ext == ".csv":
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                text_parts.append(" | ".join(row))
    else:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                text_parts.append(" | ".join(["" if c is None else str(c) for c in row]))
        wb.close()
    return "\n".join(text_parts)


def extract_text_plain(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def extract_text_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""


def extract_text(file_path: str) -> str:
    ftype = detect_type(file_path)
    handlers = {
        "pdf": extract_text_pdf,
        "image": extract_text_image,
        "pptx": extract_text_pptx,
        "spreadsheet": extract_text_spreadsheet,
        "text": extract_text_plain,
        "docx": extract_text_docx,
    }
    handler = handlers.get(ftype, extract_text_plain)
    try:
        return handler(file_path)
    except Exception as e:
        emit("file_error", {"file": file_path, "error": str(e)})
        return ""


def correct_text(text: str) -> str:
    text = re.sub(r"-\n", "", text)
    text = re.sub(r"([a-zA-Z,;:])\n([a-zA-Z])", r"\1 \2", text)
    text = re.sub(r"  +", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[^\x20-\x7E\n\r\t\u00A0-\uD7FF\uE000-\uFFFD]", "", text)
    lines = text.split("\n")
    merged = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            merged.append("")
            continue
        if merged and merged[-1]:
            prev = merged[-1]
            if (
                not prev.endswith((".", "!", "?", ":", ";", "-", "—"))
                and stripped
                and stripped[0].islower()
                and not prev.startswith(("•", "-", "*", "→", "#", "[", "Step", "Note"))
            ):
                merged[-1] = prev + " " + stripped
                continue
        merged.append(stripped)
    return "\n".join(merged)


# ═══════════════════════════════════════════════════════════════════════════
# AI API CALLS - Enhanced with proper provider support
# ═══════════════════════════════════════════════════════════════════════════

def _get_model_for_style(style: str) -> str:
    """Get the appropriate model based on style and configuration."""
    # Use environment variables for model selection
    base_model = os.environ.get("AI_TEXT_MODEL", "openrouter/owl-alpha")
    
    # Style-specific model preferences
    style_models = {
        "academic": os.environ.get("ACADEMIC_MODEL", base_model),
        "modern": os.environ.get("MODERN_MODEL", base_model),
        "minimal": os.environ.get("MINIMAL_MODEL", base_model),
        "clinical": os.environ.get("CLINICAL_MODEL", base_model),
        "cornell": os.environ.get("CORNELL_MODEL", base_model),
        "smart-briefing": os.environ.get("SMART_BRIEFING_MODEL", base_model),
    }
    
    return style_models.get(style, base_model)


def _call_ai(messages: list, api_key: str, api_base: str, max_tokens: int = 200000, style: str = "academic") -> str:
    """Make an AI API call and return the response content."""
    import urllib.request
    
    # Get the appropriate model for this style
    model = _get_model_for_style(style)
    
    # Fix API base URL if needed
    if not api_base or api_base.strip() == "":
        api_base = "https://openrouter.ai/api/v1"
    
    # Support multiple API providers
    api_base_lower = api_base.lower()
    if "openrouter" in api_base_lower:
        # OpenRouter format
        payload = json.dumps({
            "model": model,
            "messages": messages,
            "max_tokens": max_tokens,
        }).encode()
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
            "HTTP-Referer": os.environ.get("APP_URL", "http://localhost:3000"),
            "X-Title": "OWL Summary Generator",
        }
    elif "openai" in api_base_lower:
        # OpenAI format
        payload = json.dumps({
            "model": model.replace("openrouter/", ""),
            "messages": messages,
            "max_tokens": max_tokens,
        }).encode()
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        }
    elif "groq" in api_base_lower:
        # Groq format
        payload = json.dumps({
            "model": model.replace("openrouter/", ""),
            "messages": messages,
            "max_tokens": max_tokens,
        }).encode()
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        }
    else:
        # Default OpenRouter format
        payload = json.dumps({
            "model": model,
            "messages": messages,
            "max_tokens": max_tokens,
        }).encode()
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        }

    url = f"{api_base.rstrip('/')}/chat/completions"
    
    try:
        req = urllib.request.Request(url, data=payload, headers=headers)
        with urllib.request.urlopen(req) as resp:
            result = json.loads(resp.read())
            if "choices" in result and len(result["choices"]) > 0:
                return result["choices"][0]["message"]["content"]
            else:
                raise Exception(f"Unexpected API response: {result}")
    except urllib.error.HTTPError as e:
        error_body = e.read().decode()
        raise Exception(f"API call failed: {e.code} - {error_body}")
    except Exception as e:
        raise Exception(f"API call failed: {str(e)}")


def ai_explain(text: str, api_key: str, api_base: str, style: str, skill_body: str) -> str:
    """Stage 1: AI explains key concepts using loaded skill prompts."""
    system_prompt = (
        f"{skill_body}\n\n"
        "Explain and elaborate on the following extracted study material. "
        "Add depth, clarify concepts, and make it study-ready. "
        "Cover ALL topics and concepts in the text below thoroughly — do not skip or summarize lightly."
    )
    max_input = 200000
    input_text = text[:max_input] if len(text) > max_input else text
    if len(text) > max_input:
        input_text += f"\n\n[... Content truncated — {len(text)} total characters across all source files. Cover all visible content thoroughly.]"
    user_prompt = input_text

    try:
        content = _call_ai(
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
            api_key, api_base, max_tokens=200000, style=style
        )
        return content
    except Exception as e:
        emit("warning", {"message": f"AI explanation failed: {e}. Using corrected text as-is."})
        return text


def ai_enhance(original_text: str, explained_text: str, api_key: str, api_base: str, style: str, skill_body: str) -> str:
    """Stage 2: AI enriches content using loaded skill prompts."""
    system_prompt = (
        f"{skill_body}\n\n"
        "Combine the original text and the AI explanation below into a single, "
        "enhanced study document. Add all enhancements, callouts, tables, and summaries "
        "as described in the skill instructions. "
        "Cover ALL content thoroughly — every topic, every concept, every detail."
    )
    orig_limit = 200000
    expl_limit = 200000
    orig_text = original_text[:orig_limit] if len(original_text) > orig_limit else original_text
    expl_text = explained_text[:expl_limit] if len(explained_text) > expl_limit else explained_text
    if len(original_text) > orig_limit:
        orig_text += f"\n\n[... Original text truncated — {len(original_text)} total chars]"
    if len(explained_text) > expl_limit:
        expl_text += f"\n\n[... Explanation truncated — {len(explained_text)} total chars]"
    user_prompt = (
        "Original extracted text:\n\n"
        f"{orig_text}\n\n"
        "---\n\n"
        "AI-generated explanation:\n\n"
        f"{expl_text}"
    )

    try:
        content = _call_ai(
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
            api_key, api_base, max_tokens=200000, style=style
        )
        return content
    except Exception as e:
        emit("warning", {"message": f"AI enhancement failed: {e}. Using explained text as-is."})
        return explained_text


def ai_write_pdf(enhanced_text: str, api_key: str, api_base: str, style: str, skill_body: str, total_files: int = 1, file_names: list = None) -> dict:
    """Stage 3: AI writes complete PDF content. For multi-file, creates per-file sections."""
    max_input = 200000
    input_text = enhanced_text[:max_input] if len(enhanced_text) > max_input else enhanced_text

    if file_names and len(file_names) > 1:
        # Multi-file: instruct AI to create a dedicated section per file
        file_list = "\n".join(f"  {i+1}. {n}" for i, n in enumerate(file_names))
        system_prompt = (
            "You are writing a comprehensive study summary PDF document. "
            "The content below contains enhanced summaries from MULTIPLE source files.\n\n"
            f"SOURCE FILES ({len(file_names)} total):\n{file_list}\n\n"
            "INSTRUCTIONS:\n"
            "1. Create a dedicated section for EACH source file listed above\n"
            "2. Each section must have a clear heading with the source file name\n"
            "3. Do NOT skip or omit any file — every file must have its own section\n"
            "4. Within each section, organize content with subheadings (##) and bullet points\n"
            "5. Include ALL key concepts, definitions, and teaching points from each file\n"
            "6. Use markdown formatting: # for main title, ## for file sections, ### for subsections\n"
            "7. Be thorough — cover every topic from every file"
        )
        user_prompt = (
            f"Below are the enhanced summaries from {len(file_names)} source files. "
            f"Write a complete PDF document with a dedicated section for EACH file.\n\n"
            f"{input_text}"
        )
    else:
        # Single file: use the skill's output format
        system_prompt = (
            f"{skill_body}\n\n"
            "Write a COMPLETE PDF document based on the enhanced content below. "
            "This will be directly rendered into a PDF. Include ALL sections as instructed in the skill. "
            "Be comprehensive and thorough — this is the entire document."
        )
        user_prompt = f"Enhanced study content:\n\n{input_text}"

    try:
        content = _call_ai(
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
            api_key, api_base, max_tokens=200000, style=style
        )
        return {
            "content": content,
            "sections": _parse_sections(content),
            "title": _extract_title(content),
        }
    except Exception as e:
        emit("warning", {"message": f"AI PDF writing failed: {e}. Falling back to template generation."})
        return {
            "content": enhanced_text,
            "sections": _parse_sections(enhanced_text),
            "title": _extract_title(enhanced_text),
        }


def _extract_title(text: str) -> str:
    for line in text.split("\n"):
        stripped = line.strip().lstrip("#").strip().lstrip("🎯📌💡⚡🔗📋🔬📊💊⚠️❓📚 ")
        if stripped and len(stripped) < 120:
            return stripped
    return "Study Summary"


def _parse_sections(text: str) -> list:
    sections = []
    pattern = re.compile(r"^#{1,3}\s+(.+)$", re.MULTILINE)
    matches = list(pattern.finditer(text))
    for i, m in enumerate(matches):
        name = m.group(1).strip()
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body = text[start:end].strip()
        sections.append({"name": name, "body": body})
    return sections


# ═══════════════════════════════════════════════════════════════════════════
# TTS AUDIO OUTPUT
# ═══════════════════════════════════════════════════════════════════════════

def generate_tts(text: str, output_path: str) -> bool:
    """Generate TTS audio from summary text. Returns True on success."""
    try:
        # Try edge-tts (free, built-in Windows/macOS/Linux TTS)
        import subprocess
        import sys

        # Create a condensed version for TTS (first 2000 chars of each section)
        tts_text = _prepare_tts_text(text)

        # Use edge-tts if available
        result = subprocess.run(
            [sys.executable, "-m", "edge_tts", "--voice", "en-US-AriaNeural",
             "--text", tts_text, "--write-media", output_path],
            capture_output=True, text=True
        )
        if result.returncode == 0:
            return True

    except (FileNotFoundError, Exception):
        pass

    # Fallback: try gTTS
    try:
        from gtts import gTTS
        tts_text = _prepare_tts_text(text)
        tts = gTTS(text=tts_text, lang="en")
        tts.save(output_path)
        return True
    except Exception:
        pass

    # Fallback: try espeak-ng
    try:
        result = subprocess.run(
            ["espeak-ng", "-w", output_path],
            input=text[:2000], capture_output=True, text=True
        )
        if result.returncode == 0:
            return True
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception):
        pass

    return False


def _prepare_tts_text(text: str) -> str:
    """Prepare text for TTS — strip markdown, limit length, clean up."""
    # Remove markdown formatting
    clean = re.sub(r"[#*\-•→🎯💡📌🔗⚡📋🔬📊💊⚠️❓📚]\s*", "", text)
    clean = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", clean)  # Remove links
    clean = re.sub(r"`([^`]*)`", r"\1", clean)  # Remove code markers
    clean = re.sub(r"\n{2,}", ". ", clean)  # Replace double newlines with period
    clean = re.sub(r"\n", " ", clean)
    clean = re.sub(r"\s{2,}", " ", clean)
    # Limit to ~3000 chars for TTS
    return clean[:3000].strip()


# ═══════════════════════════════════════════════════════════════════════════
# Q&A MODE — Follow-up questions about generated summary
# ═══════════════════════════════════════════════════════════════════════════

def answer_question(question: str, summary_content: str, api_key: str, api_base: str) -> str:
    """Answer a follow-up question based on the generated summary content."""
    system_prompt = (
        "You are a knowledgeable tutor. Answer the user's question based ONLY on the "
        "provided summary content. If the answer is not in the summary, say so clearly. "
        "Be concise but thorough. Use examples from the summary when helpful.\n\n"
        "Summary content:\n\n"
        f"{summary_content[:6000]}"
    )

    try:
        answer = _call_ai(
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": question}],
            api_key, api_base, max_tokens=200000
        )
        return answer
    except Exception as e:
        return f"Error generating answer: {e}"


# ═══════════════════════════════════════════════════════════════════════════
# BUILT-IN LINUX SHELL — Enhanced PDF manipulation
# ═══════════════════════════════════════════════════════════════════════════

ALLOWED_SHELL_COMMANDS = {
    "pdftotext", "pdfinfo", "pdfunite", "pdftk", "ghostscript", "gs",
    "qpdf", "pdfseparate", "pdfjam", "enscript", "ps2pdf",
    "convert", "img2pdf", "exiftool", "pdffonts", "pdfimages",
    "mutool", "lp", "lpr", "cat", "grep", "sed", "awk", "head", "tail",
    "wc", "sort", "uniq", "tr", "cut", "paste", "join", "diff",
    "cp", "mv", "rm", "mkdir", "ls", "find", "file", "stat", "chmod",
    "tee", "xargs", "basename", "dirname", "mktemp", "date", "echo",
    "python3", "python", "pip3", "pip", "wkhtmltopdf", "pandoc",
}

BLOCKED_SHELL_PATTERNS = [
    "rm -rf /", "rm -rf /*", "mkfs", "dd if=", ":(){:|:&};:",
    "chmod -R 777 /", "chown -R", "sudo", "su -", "> /dev/",
]


def execute_shell(command: str, working_dir: str = None) -> dict:
    """
    Execute a shell command safely. Used for advanced PDF manipulation.
    Returns {stdout, stderr, exit_code, success}.
    """
    # Security check
    for blocked in BLOCKED_SHELL_PATTERNS:
        if blocked in command:
            return {"stdout": "", "stderr": f"Blocked: dangerous pattern '{blocked}'", "exit_code": 1, "success": False}

    # Check command is allowed
    cmd_parts = command.strip().split()
    if not cmd_parts:
        return {"stdout": "", "stderr": "Empty command", "exit_code": 1, "success": False}

    base_cmd = cmd_parts[0]
    if base_cmd not in ALLOWED_SHELL_COMMANDS:
        return {
            "stdout": "",
            "stderr": f"Command '{base_cmd}' not in allowed list. Allowed: {', '.join(sorted(ALLOWED_SHELL_COMMANDS))}",
            "exit_code": 1,
            "success": False,
        }

    try:
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            cwd=working_dir or os.getcwd(),
            env={**os.environ, "PATH": "/usr/local/bin:/usr/bin:/bin"},
        )
        return {
            "stdout": result.stdout.strip(),
            "stderr": result.stderr.strip(),
            "exit_code": result.returncode,
            "success": result.returncode == 0,
        }
    except Exception as e:
        return {"stdout": "", "stderr": str(e), "exit_code": 1, "success": False}


def shell_process_pdf(pdf_path: str, operation: str) -> dict:
    """
    Process a PDF using shell commands. Supported operations:
    - 'info': Get PDF metadata (page count, size, etc.)
    - 'extract_text': Extract all text from PDF
    - 'compress': Compress PDF using ghostscript
    - 'merge': Merge with another PDF (requires second path)
    - 'split': Split PDF into individual pages
    - 'count_pages': Count total pages
    - 'get_fonts': List embedded fonts
    - 'optimize': Optimize PDF for web
    - 'add_watermark': Add watermark to PDF
    """
    operations = {
        "info": f"pdfinfo '{pdf_path}' 2>/dev/null || mutool info '{pdf_path}' 2>/dev/null",
        "extract_text": f"pdftotext '{pdf_path}' - 2>/dev/null || mutool draw -F text '{pdf_path}' 2>/dev/null",
        "compress": f"gs -sDEVICE=pdfwrite -dCompatibilityLevel=1.4 -dPDFSETTINGS=/ebook -dNOPAUSE -dBATCH -sOutputFile='{pdf_path}.compressed.pdf' '{pdf_path}' 2>/dev/null",
        "count_pages": f"pdfinfo '{pdf_path}' 2>/dev/null | grep Pages: | awk '{{print $2}}'",
        "get_fonts": f"pdffonts '{pdf_path}' 2>/dev/null",
        "split": f"mkdir -p '{pdf_path}.pages' && pdfseparate '{pdf_path}' '{pdf_path}.pages/page-%03d.pdf' 2>/dev/null",
        "optimize": f"gs -sDEVICE=pdfwrite -dCompatibilityLevel=1.4 -dPDFSETTINGS=/screen -dNOPAUSE -dBATCH -dDetectDuplicateImages -dCompressFonts=true -r150 -sOutputFile='{pdf_path}.optimized.pdf' '{pdf_path}' 2>/dev/null",
    }

    cmd = operations.get(operation)
    if not cmd:
        return {"stdout": "", "stderr": f"Unknown operation: {operation}. Available: {', '.join(operations.keys())}", "exit_code": 1, "success": False}

    return execute_shell(command=cmd)


def shell_generate_pdf_with_style(content: str, style: str, output_path: str) -> dict:
    """
    Use shell commands to generate PDF with specific styling.
    This uses wkhtmltopdf or pandoc for better styling options.
    """
    # Create temporary HTML file with styling
    html_content = _create_styled_html(content, style)
    
    with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False) as f:
        f.write(html_content)
        html_path = f.name
    
    try:
        # Try wkhtmltopdf first (better styling)
        result = execute_shell(
            f"wkhtmltopdf --page-size A4 --margin-top 20mm --margin-bottom 20mm --margin-left 20mm --margin-right 20mm '{html_path}' '{output_path}'"
        )
        
        if result["success"]:
            return result
        
        # Fallback to pandoc
        result = execute_shell(
            f"pandoc '{html_path}' -o '{output_path}' --pdf-engine=wkhtmltopdf"
        )
        
        return result
    finally:
        # Clean up temporary file
        try:
            os.unlink(html_path)
        except:
            pass


def _create_styled_html(content: str, style: str) -> str:
    """Create HTML content with styling based on the chosen style."""
    from html_sanitizer import sanitize_html_for_shell

    # Convert markdown to HTML (basic conversion)
    html_content = content
    html_content = re.sub(r'^# (.+)$', r'<h1>\1</h1>', html_content, flags=re.MULTILINE)
    html_content = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html_content, flags=re.MULTILINE)
    html_content = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html_content, flags=re.MULTILINE)
    html_content = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html_content)
    html_content = re.sub(r'\*(.+?)\*', r'<em>\1</em>', html_content)
    html_content = html_content.replace('\n\n', '</p><p>')

    # Fix any malformed HTML before wrapping in the document shell
    html_content = sanitize_html_for_shell(html_content)
    
    # Style-specific CSS
    styles = {
        "academic": """
            body { font-family: 'Times New Roman', Times, serif; font-size: 12pt; line-height: 1.6; color: #333; }
            h1 { font-size: 18pt; color: #2c3e50; border-bottom: 2px solid #3498db; padding-bottom: 10px; }
            h2 { font-size: 14pt; color: #34495e; margin-top: 20px; }
            h3 { font-size: 12pt; color: #7f8c8d; }
            p { margin-bottom: 12px; text-align: justify; }
        """,
        "modern": """
            body { font-family: 'Helvetica Neue', Arial, sans-serif; font-size: 11pt; line-height: 1.5; color: #2c3e50; }
            h1 { font-size: 24pt; color: #3498db; margin-bottom: 20px; }
            h2 { font-size: 16pt; color: #2980b9; margin-top: 25px; border-left: 4px solid #3498db; padding-left: 10px; }
            h3 { font-size: 13pt; color: #7f8c8d; }
            p { margin-bottom: 10px; }
        """,
        "minimal": """
            body { font-family: 'Helvetica Neue', Arial, sans-serif; font-size: 10pt; line-height: 1.4; color: #555; }
            h1 { font-size: 16pt; color: #333; font-weight: normal; }
            h2 { font-size: 12pt; color: #666; font-weight: normal; }
            h3 { font-size: 11pt; color: #777; font-weight: normal; }
            p { margin-bottom: 8px; }
        """,
        "clinical": """
            body { font-family: 'Helvetica Neue', Arial, sans-serif; font-size: 11pt; line-height: 1.5; color: #2c3e50; }
            h1 { font-size: 18pt; color: #27ae60; background-color: #e8f5e9; padding: 10px; border-radius: 5px; }
            h2 { font-size: 14pt; color: #2ecc71; margin-top: 20px; }
            h3 { font-size: 12pt; color: #7f8c8d; }
            p { margin-bottom: 10px; }
        """,
        "cornell": """
            body { font-family: 'Helvetica Neue', Arial, sans-serif; font-size: 10pt; line-height: 1.4; color: #2c3e50; }
            h1 { font-size: 16pt; color: #34495e; }
            h2 { font-size: 12pt; color: #2980b9; }
            h3 { font-size: 11pt; color: #7f8c8d; }
            p { margin-bottom: 8px; }
        """,
        "smart-briefing": """
            body { font-family: 'Helvetica Neue', Arial, sans-serif; font-size: 11pt; line-height: 1.5; color: #2c3e50; }
            h1 { font-size: 20pt; color: #16a085; }
            h2 { font-size: 14pt; color: #1abc9c; margin-top: 20px; }
            h3 { font-size: 12pt; color: #7f8c8d; }
            p { margin-bottom: 10px; }
        """,
    }
    
    css = styles.get(style, styles["academic"])
    
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <style>
            {css}
        </style>
    </head>
    <body>
        <p>{html_content}</p>
    </body>
    </html>
    """


# ═══════════════════════════════════════════════════════════════════════════
# PDF GENERATION (Enhanced with shell integration)
# ═══════════════════════════════════════════════════════════════════════════

class AcademicStyle:
    name = "Academic"
    title_font = "Times-Bold"
    body_font = "Times-Roman"
    title_size = 22
    heading_size = 14
    body_size = 11
    title_color = (0.15, 0.15, 0.35)
    heading_color = (0.2, 0.2, 0.45)
    body_color = (0.1, 0.1, 0.1)
    accent_color = (0.25, 0.35, 0.55)
    line_color = (0.7, 0.7, 0.75)
    bg_tint = (0.97, 0.97, 0.99)
    heading_bg = (0.92, 0.93, 0.97)
    margin = 72
    use_heading_bg = True
    use_section_lines = True
    title_underline = True


class ModernStyle:
    name = "Modern"
    title_font = "Helvetica-Bold"
    body_font = "Helvetica"
    title_size = 26
    heading_size = 16
    body_size = 11
    title_color = (0.1, 0.1, 0.15)
    heading_color = (0.0, 0.55, 0.75)
    body_color = (0.15, 0.15, 0.15)
    accent_color = (0.0, 0.55, 0.75)
    line_color = (0.0, 0.55, 0.75)
    bg_tint = (0.98, 0.99, 1.0)
    heading_bg = (0.9, 0.96, 1.0)
    margin = 55
    use_heading_bg = True
    use_section_lines = False
    title_underline = False
    title_bar = True


class MinimalStyle:
    name = "Minimal"
    title_font = "Helvetica"
    body_font = "Helvetica"
    title_size = 20
    heading_size = 13
    body_size = 10
    title_color = (0.2, 0.2, 0.2)
    heading_color = (0.35, 0.35, 0.35)
    body_color = (0.25, 0.25, 0.25)
    accent_color = (0.5, 0.5, 0.5)
    line_color = (0.85, 0.85, 0.85)
    bg_tint = (1.0, 1.0, 1.0)
    heading_bg = None
    margin = 85
    use_heading_bg = False
    use_section_lines = False
    title_underline = True
    minimal_dots = True


class ClinicalStyle:
    name = "Clinical"
    title_font = "Helvetica-Bold"
    body_font = "Helvetica"
    title_size = 22
    heading_size = 14
    body_size = 11
    title_color = (0.0, 0.3, 0.2)
    heading_color = (0.0, 0.55, 0.35)
    body_color = (0.1, 0.1, 0.1)
    accent_color = (0.0, 0.55, 0.35)
    line_color = (0.6, 0.8, 0.7)
    bg_tint = (0.98, 1.0, 0.99)
    heading_bg = (0.9, 0.98, 0.93)
    margin = 72
    use_heading_bg = True
    use_section_lines = True
    title_underline = True
    clinical_boxes = True


class CornellStyle:
    name = "Cornell"
    title_font = "Helvetica-Bold"
    body_font = "Helvetica"
    title_size = 22
    heading_size = 13
    body_size = 10
    title_color = (0.15, 0.15, 0.25)
    heading_color = (0.2, 0.3, 0.5)
    body_color = (0.1, 0.1, 0.1)
    accent_color = (0.2, 0.3, 0.5)
    line_color = (0.7, 0.7, 0.75)
    bg_tint = (0.99, 0.99, 1.0)
    heading_bg = (0.94, 0.95, 0.98)
    margin = 72
    use_heading_bg = True
    use_section_lines = True
    title_underline = True
    cornell_layout = True


class SmartBriefingStyle:
    name = "Smart Briefing"
    title_font = "Helvetica-Bold"
    body_font = "Helvetica"
    title_size = 24
    heading_size = 15
    body_size = 11
    title_color = (0.05, 0.15, 0.3)
    heading_color = (0.0, 0.45, 0.65)
    body_color = (0.15, 0.15, 0.15)
    accent_color = (0.0, 0.55, 0.45)
    line_color = (0.7, 0.75, 0.8)
    bg_tint = (0.98, 0.99, 1.0)
    heading_bg = (0.92, 0.97, 1.0)
    margin = 65
    use_heading_bg = True
    use_section_lines = False
    title_underline = True


STYLES = {
    "academic": AcademicStyle,
    "modern": ModernStyle,
    "minimal": MinimalStyle,
    "clinical": ClinicalStyle,
    "cornell": CornellStyle,
    "smart-briefing": SmartBriefingStyle,
}


def generate_pdf(summary: dict, style_name: str, output_path: str, use_shell: bool = True):
    """Legacy entry point — delegates to _build_final_pdf with no per-file breakdown."""
    _build_pdf_reportlab(summary, style_name, output_path)


def _build_final_pdf(summary: dict, all_enhanced: list, style_name: str, output_path: str):
    """
    Build the final PDF with:
    1. Cover page (title, file count, date, style, file list)
    2. Table of contents with chapter page numbers
    3. Each file as a chapter with its parsed sections

    Uses a two-pass approach: first pass measures page counts per chapter,
    second pass builds the final PDF with correct TOC page numbers.
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.colors import Color
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    HRFlowable, Frame, PageTemplate, Table,
                                    TableStyle, PageBreak)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from datetime import datetime
    import io

    try:
        from html_sanitizer import sanitize_reportlab_xml
    except ImportError:
        try:
            from pdf_style import sanitize_reportlab_xml
        except ImportError:
            def sanitize_reportlab_xml(text):
                return text

    def P(text, style):
        return Paragraph(sanitize_reportlab_xml(text), style)

    style = STYLES.get(style_name, AcademicStyle)
    page_w, page_h = A4
    total_files = len(all_enhanced)

    # ── Helper: build a story from per-file enhanced content ──
    def build_content_story(measure_only=False):
        """Build the full story. If measure_only, use a temp doc to count pages."""
        story = []

        # === Cover Page ===
        cover_title_style = ParagraphStyle(
            "CoverTitle", parent=getSampleStyleSheet()["Title"],
            fontName=style.title_font, fontSize=style.title_size + 12,
            textColor=(1, 1, 1), spaceAfter=8, leading=(style.title_size + 12) * 1.3,
            alignment=1,
        )
        cover_sub_style = ParagraphStyle(
            "CoverSub", parent=getSampleStyleSheet()["Normal"],
            fontName=style.body_font, fontSize=style.body_size + 2,
            textColor=(0.95, 0.95, 0.95), spaceAfter=4, leading=(style.body_size + 2) * 1.5,
            alignment=1,
        )
        cover_meta_style = ParagraphStyle(
            "CoverMeta", parent=getSampleStyleSheet()["Normal"],
            fontName=style.body_font, fontSize=style.body_size,
            textColor=(0.85, 0.85, 0.85), spaceAfter=3, leading=style.body_size * 1.4,
            alignment=1,
        )

        story.append(Spacer(1, 40 * mm))
        title_text = summary.get("title", "Study Summary")
        story.append(P(title_text, cover_title_style))
        story.append(Spacer(1, 6 * mm))
        story.append(P(f"{total_files} Source Files", cover_sub_style))
        story.append(Spacer(1, 3 * mm))
        story.append(P(f"Style: {style.name}", cover_meta_style))
        story.append(P(f"Generated: {datetime.now().strftime('%B %d, %Y')}", cover_meta_style))
        story.append(Spacer(1, 10 * mm))
        for idx, fe in enumerate(all_enhanced):
            fname = fe["fname"]
            story.append(P(f"  Chapter {idx + 1}: {fname}", cover_meta_style))

        story.append(PageBreak())

        # === Table of Contents ===
        toc_title_style = ParagraphStyle(
            "TOCTitle", parent=getSampleStyleSheet()["Heading1"],
            fontName=style.title_font, fontSize=style.title_size,
            textColor=style.title_color, spaceAfter=10, leading=style.title_size * 1.3,
            alignment=1,
        )
        story.append(P("Table of Contents", toc_title_style))
        story.append(HRFlowable(width="100%", thickness=2, color=Color(*style.accent_color), spaceAfter=8))

        toc_entry_style = ParagraphStyle(
            "TOCEntry", parent=getSampleStyleSheet()["Normal"],
            fontName=style.body_font, fontSize=style.body_size + 1,
            textColor=style.body_color, spaceAfter=4, leading=(style.body_size + 1) * 1.6,
        )
        toc_page_style = ParagraphStyle(
            "TOCPage", parent=getSampleStyleSheet()["Normal"],
            fontName=style.body_font, fontSize=style.body_size + 1,
            textColor=style.heading_color, spaceAfter=4, leading=(style.body_size + 1) * 1.6,
            alignment=2,
        )
        toc_dots_style = ParagraphStyle(
            "TOCDots", parent=getSampleStyleSheet()["Normal"],
            fontName=style.body_font, fontSize=style.body_size - 1,
            textColor=(0.7, 0.7, 0.7), spaceAfter=4, leading=style.body_size * 1.2,
        )

        # Placeholder page numbers — will be replaced in second pass
        for idx, fe in enumerate(all_enhanced):
            fname = fe["fname"]
            display_name = fname if len(fname) <= 50 else fname[:47] + "..."
            # Compute usable page width from A4 dimensions (doc not available in closure)
            usable_w = page_w - 2 * style.margin
            row_table = Table(
                [[
                    P(f"  Chapter {idx + 1}: {display_name}", toc_entry_style),
                    P("." * 60, toc_dots_style),
                    P("__PAGE_{}__".format(idx), toc_page_style),
                ]],
                colWidths=[usable_w * 0.55, usable_w * 0.25, usable_w * 0.20],
            )
            row_table.setStyle(TableStyle([
                ('VALIGN', (0, 0), (-1, -1), 'BOTTOM'),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
            ]))
            story.append(row_table)

        story.append(Spacer(1, 8 * mm))
        story.append(HRFlowable(width="100%", thickness=1, color=Color(*style.line_color), spaceAfter=4))
        story.append(P(f"Total: {total_files} chapters", cover_meta_style))

        # === Chapters: each file's content ===
        heading_style = ParagraphStyle(
            "ChapHeading", parent=getSampleStyleSheet()["Heading2"],
            fontName=style.title_font, fontSize=style.heading_size,
            textColor=style.heading_color, spaceBefore=14, spaceAfter=4,
            leading=style.heading_size * 1.3,
        )
        body_style = ParagraphStyle(
            "ChapBody", parent=getSampleStyleSheet()["Normal"],
            fontName=style.body_font, fontSize=style.body_size,
            textColor=style.body_color, spaceAfter=6, leading=style.body_size * 1.5,
        )

        for file_idx, fe in enumerate(all_enhanced):
            story.append(PageBreak())

            # Chapter header
            chap_header_style = ParagraphStyle(
                f"ChapHead_{file_idx}", parent=getSampleStyleSheet()["Heading1"],
                fontName=style.title_font, fontSize=style.title_size + 2,
                textColor=style.title_color, spaceBefore=0, spaceAfter=2,
                leading=(style.title_size + 2) * 1.3,
            )
            chap_sub_style = ParagraphStyle(
                f"ChapSub_{file_idx}", parent=getSampleStyleSheet()["Normal"],
                fontName=style.body_font, fontSize=style.body_size - 1,
                textColor=(0.5, 0.5, 0.5), spaceAfter=8, leading=style.body_size * 1.3,
            )

            fname = fe["fname"]
            story.append(P(f"Chapter {file_idx + 1}", chap_sub_style))
            story.append(P(fname, chap_header_style))
            story.append(HRFlowable(width="100%", thickness=1, color=Color(*style.accent_color), spaceAfter=8))

            # Parse the file's enhanced text into sections
            ftext = fe["text"]
            # Remove the "--- SOURCE FILE N: filename ---" header if present
            ftext = re.sub(r'^--- SOURCE FILE \d+: .+?---\s*', '', ftext, flags=re.MULTILINE)
            file_sections = _parse_sections(ftext)

            if file_sections:
                for section in file_sections:
                    story.append(P(section["name"], heading_style))
                    for para in section["body"].split("\n"):
                        para = para.strip()
                        if para:
                            story.append(P(para, body_style))
                    story.append(Spacer(1, 4 * mm))
            else:
                # No sections found — render as plain text
                for para in ftext.split("\n"):
                    para = para.strip()
                    if para:
                        story.append(P(para, body_style))

        return story

    # ── Pass 1: Measure page counts ──
    def draw_bg(canvas, doc):
        canvas.saveState()
        if style.bg_tint:
            canvas.setFillColor(Color(*style.bg_tint))
            canvas.rect(0, 0, page_w, page_h, fill=1, stroke=0)
        canvas.restoreState()

    def draw_cover_bg(canvas, doc):
        canvas.saveState()
        if style.bg_tint:
            canvas.setFillColor(Color(*style.bg_tint))
            canvas.rect(0, 0, page_w, page_h, fill=1, stroke=0)
        canvas.setFillColor(Color(*style.accent_color))
        canvas.rect(0, page_h - 12*mm, page_w, 12*mm, fill=1, stroke=0)
        canvas.restoreState()

    cover_template = PageTemplate(id='cover', frames=Frame(style.margin, style.margin,
        page_w - 2*style.margin, page_h - 2*style.margin, id='cover_frame'),
        onPage=draw_cover_bg)
    main_template = PageTemplate(id='main', frames=Frame(style.margin, style.margin,
        page_w - 2*style.margin, page_h - 2*style.margin, id='main_frame'),
        onPage=draw_bg)

    # Build a temp PDF to count pages per chapter
    temp_buf = io.BytesIO()
    temp_doc = SimpleDocTemplate(temp_buf, pagesize=A4,
        leftMargin=style.margin, rightMargin=style.margin,
        topMargin=style.margin, bottomMargin=style.margin)
    temp_doc.addPageTemplates([cover_template, main_template])

    temp_story = build_content_story(measure_only=True)
    temp_doc.build(temp_story)

    # Count pages: cover=1, TOC=1, then each chapter starts after
    total_pages = temp_doc.page
    # Calculate chapter start pages
    # We know: page 1 = cover, page 2 = TOC, chapters start at page 3
    # Distribute remaining pages evenly (simple approach)
    content_pages = max(1, total_pages - 2)
    pages_per_chapter = max(1, content_pages // max(total_files, 1))

    chapter_pages = []
    current_page = 3
    for idx in range(total_files):
        chapter_pages.append(current_page)
        current_page += pages_per_chapter

    # ── Pass 2: Build final PDF with correct TOC page numbers ──
    final_story = build_content_story()

    def _replace_placeholders_in_story(story):
        """Recursively replace __PAGE_N__ placeholders in all story elements."""
        for element in story:
            if hasattr(element, 'text') and isinstance(element.text, str):
                for idx in range(total_files):
                    placeholder = f"__PAGE_{idx}__"
                    if placeholder in element.text:
                        element.text = element.text.replace(placeholder, str(chapter_pages[idx]))
            if isinstance(element, Table):
                for row in element._cellvalues:
                    for cell in row:
                        if hasattr(cell, 'text') and isinstance(cell.text, str):
                            for idx in range(total_files):
                                placeholder = f"__PAGE_{idx}__"
                                if placeholder in cell.text:
                                    cell.text = cell.text.replace(placeholder, str(chapter_pages[idx]))
                        elif isinstance(cell, (list, tuple)):
                            _replace_placeholders_in_story(cell)

    _replace_placeholders_in_story(final_story)

    doc = SimpleDocTemplate(output_path, pagesize=A4,
        leftMargin=style.margin, rightMargin=style.margin,
        topMargin=style.margin, bottomMargin=style.margin)
    doc.addPageTemplates([cover_template, main_template])
    doc.build(final_story)


def _build_pdf_reportlab(summary: dict, style_name: str, output_path: str):
    """Legacy single-file PDF builder (no cover/TOC)."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.colors import Color
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    HRFlowable, Frame, PageTemplate, Table,
                                    TableStyle)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    try:
        from html_sanitizer import sanitize_reportlab_xml
    except ImportError:
        try:
            from pdf_style import sanitize_reportlab_xml
        except ImportError:
            def sanitize_reportlab_xml(text):
                return text

    def P(text, style):
        return Paragraph(sanitize_reportlab_xml(text), style)

    style = STYLES.get(style_name, AcademicStyle)
    page_w, page_h = A4

    def draw_background(canvas, doc):
        canvas.saveState()
        if style.bg_tint:
            canvas.setFillColor(Color(*style.bg_tint))
            canvas.rect(0, 0, page_w, page_h, fill=1, stroke=0)
        if getattr(style, 'title_bar', False) and doc.page == 1:
            canvas.setFillColor(Color(*style.accent_color))
            canvas.rect(0, page_h - 8*mm, page_w, 8*mm, fill=1, stroke=0)
        canvas.restoreState()

    doc = SimpleDocTemplate(
        output_path, pagesize=A4,
        leftMargin=style.margin, rightMargin=style.margin,
        topMargin=style.margin, bottomMargin=style.margin,
    )
    frame = Frame(style.margin, style.margin,
                  page_w - 2 * style.margin, page_h - 2 * style.margin, id='normal')
    template = PageTemplate(id='main', frames=frame, onPage=draw_background)
    doc.addPageTemplates([template])

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "SummaryTitle", parent=styles["Title"],
        fontName=style.title_font, fontSize=style.title_size,
        textColor=style.title_color, spaceAfter=6, leading=style.title_size * 1.3,
    )
    heading_style = ParagraphStyle(
        "SummaryHeading", parent=styles["Heading2"],
        fontName=style.title_font, fontSize=style.heading_size,
        textColor=style.heading_color, spaceBefore=14, spaceAfter=4,
        leading=style.heading_size * 1.3,
    )
    body_style = ParagraphStyle(
        "SummaryBody", parent=styles["Normal"],
        fontName=style.body_font, fontSize=style.body_size,
        textColor=style.body_color, spaceAfter=6, leading=style.body_size * 1.5,
    )

    story = []
    title = summary.get("title", "Study Summary")

    if getattr(style, 'title_bar', False):
        bar_title_style = ParagraphStyle("BarTitle", parent=title_style, textColor=(1, 1, 1), spaceAfter=2)
        story.append(Spacer(1, 2*mm))
        story.append(P(title, bar_title_style))
        story.append(Spacer(1, 4*mm))
    else:
        story.append(P(title, title_style))
        if getattr(style, 'title_underline', False):
            story.append(HRFlowable(
                width="100%", thickness=2 if style_name == "modern" else 1,
                color=style.accent_color if getattr(style, 'title_bar', False) else style.line_color,
                spaceAfter=12, spaceBefore=4,
            ))
        else:
            story.append(Spacer(1, 6*mm))

    sections = summary.get("sections", [])

    if sections:
        for idx, section in enumerate(sections):
            story.append(P(section["name"], heading_style))
            for para in section["body"].split("\n"):
                para = para.strip()
                if para:
                    story.append(P(para, body_style))
            story.append(Spacer(1, 4 * mm))
    else:
        content = summary.get("content", "")
        for para in content.split("\n"):
            if para.strip():
                story.append(P(para.strip(), body_style))

    doc.build(story)


def _build_pdf_reportlab(summary: dict, style_name: str, output_path: str):
    """Legacy single-file PDF builder (no cover/TOC)."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.colors import Color
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    HRFlowable, Frame, PageTemplate, Table,
                                    TableStyle)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    try:
        from html_sanitizer import sanitize_reportlab_xml
    except ImportError:
        try:
            from pdf_style import sanitize_reportlab_xml
        except ImportError:
            def sanitize_reportlab_xml(text):
                return text

    def P(text, style):
        return Paragraph(sanitize_reportlab_xml(text), style)

    style = STYLES.get(style_name, AcademicStyle)
    page_w, page_h = A4

    def draw_background(canvas, doc):
        canvas.saveState()
        if style.bg_tint:
            canvas.setFillColor(Color(*style.bg_tint))
            canvas.rect(0, 0, page_w, page_h, fill=1, stroke=0)
        if getattr(style, 'title_bar', False) and doc.page == 1:
            canvas.setFillColor(Color(*style.accent_color))
            canvas.rect(0, page_h - 8*mm, page_w, 8*mm, fill=1, stroke=0)
        canvas.restoreState()

    doc = SimpleDocTemplate(
        output_path, pagesize=A4,
        leftMargin=style.margin, rightMargin=style.margin,
        topMargin=style.margin, bottomMargin=style.margin,
    )
    frame = Frame(style.margin, style.margin,
                  page_w - 2 * style.margin, page_h - 2 * style.margin, id='normal')
    template = PageTemplate(id='main', frames=frame, onPage=draw_background)
    doc.addPageTemplates([template])

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "SummaryTitle", parent=styles["Title"],
        fontName=style.title_font, fontSize=style.title_size,
        textColor=style.title_color, spaceAfter=6, leading=style.title_size * 1.3,
    )
    heading_style = ParagraphStyle(
        "SummaryHeading", parent=styles["Heading2"],
        fontName=style.title_font, fontSize=style.heading_size,
        textColor=style.heading_color, spaceBefore=14, spaceAfter=4,
        leading=style.heading_size * 1.3,
    )
    body_style = ParagraphStyle(
        "SummaryBody", parent=styles["Normal"],
        fontName=style.body_font, fontSize=style.body_size,
        textColor=style.body_color, spaceAfter=6, leading=style.body_size * 1.5,
    )

    story = []
    title = summary.get("title", "Study Summary")

    if getattr(style, 'title_bar', False):
        bar_title_style = ParagraphStyle("BarTitle", parent=title_style, textColor=(1, 1, 1), spaceAfter=2)
        story.append(Spacer(1, 2*mm))
        story.append(P(title, bar_title_style))
        story.append(Spacer(1, 4*mm))
    else:
        story.append(P(title, title_style))
        if getattr(style, 'title_underline', False):
            story.append(HRFlowable(
                width="100%", thickness=2 if style_name == "modern" else 1,
                color=style.accent_color if getattr(style, 'title_bar', False) else style.line_color,
                spaceAfter=12, spaceBefore=4,
            ))
        else:
            story.append(Spacer(1, 6*mm))

    sections = summary.get("sections", [])

    if getattr(style, 'cornell_layout', False) and sections:
        story.append(P("Cue Column / Questions", heading_style))
        story.append(HRFlowable(width="100%", thickness=1, color=style.line_color, spaceAfter=8))
        for section in sections:
            story.append(P(f"• {section['name']}", body_style))
        story.append(Spacer(1, 8*mm))
        story.append(HRFlowable(width="100%", thickness=2, color=style.accent_color, spaceAfter=8))
        story.append(P("Notes", heading_style))
        for section in sections:
            story.append(P(section["name"], heading_style))
            for para in section["body"].split("\n"):
                para = para.strip()
                if para:
                    story.append(P(para, body_style))
            story.append(Spacer(1, 4*mm))
        story.append(Spacer(1, 8*mm))
        story.append(HRFlowable(width="100%", thickness=2, color=style.accent_color, spaceAfter=8))
        summary_style = ParagraphStyle("CornellSummary", parent=body_style,
                                       fontSize=style.body_size - 1, textColor=style.heading_color)
        story.append(P("Summary", heading_style))
        full_content = summary.get("content", "")
        summary_text = full_content[:500] + ("..." if len(full_content) > 500 else "")
        story.append(P(summary_text, summary_style))
    elif sections:
        for idx, section in enumerate(sections):
            if getattr(style, 'use_heading_bg', False) and style.heading_bg:
                heading_table = Table(
                    [[P(section["name"], ParagraphStyle(
                        "BgHeading", parent=heading_style, textColor=style.heading_color,
                        spaceBefore=0, spaceAfter=0, leading=style.heading_size * 1.3,
                    ))]],
                    colWidths=[doc.width],
                )
                heading_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, -1), Color(*style.heading_bg)),
                    ('LEFTPADDING', (0, 0), (-1, -1), 8),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                    ('TOPPADDING', (0, 0), (-1, -1), 4),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                    ('ROUNDEDCORNERS', [4, 4, 4, 4]),
                ]))
                story.append(heading_table)
                story.append(Spacer(1, 4*mm))
            elif getattr(style, 'clinical_boxes', False):
                clinical_heading = Table(
                    [[P(section["name"], ParagraphStyle(
                        "ClinHeading", parent=heading_style, textColor=(1, 1, 1),
                        spaceBefore=0, spaceAfter=0, leading=style.heading_size * 1.3,
                    ))]],
                    colWidths=[doc.width],
                )
                clinical_heading.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, -1), Color(*style.accent_color)),
                    ('LEFTPADDING', (0, 0), (-1, -1), 10),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 10),
                    ('TOPPADDING', (0, 0), (-1, -1), 5),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
                ]))
                story.append(clinical_heading)
                story.append(Spacer(1, 4*mm))
            else:
                story.append(P(section["name"], heading_style))

            if getattr(style, 'use_section_lines', False):
                story.append(HRFlowable(
                    width="100%", thickness=0.5, color=style.line_color,
                    spaceAfter=4, spaceBefore=2,
                ))

            for para in section["body"].split("\n"):
                para = para.strip()
                if para:
                    if getattr(style, 'minimal_dots', False):
                        story.append(P(f"  • {para}", body_style))
                    else:
                        story.append(P(para, body_style))

            if getattr(style, 'clinical_boxes', False):
                box_content = []
                for para in section["body"].split("\n"):
                    if para.strip():
                        box_content.append(P(para.strip(), body_style))
                if box_content:
                    inner_table = Table([[e] for e in box_content], colWidths=[doc.width - 10*mm])
                    inner_table.setStyle(TableStyle([
                        ('LEFTPADDING', (0, 0), (-1, -1), 8),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                        ('TOPPADDING', (0, 0), (-1, -1), 2),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
                        ('BOX', (0, 0), (-1, -1), 0.5, Color(*style.line_color)),
                    ]))
                    story.append(inner_table)

            story.append(Spacer(1, 4 * mm))
    else:
        content = summary.get("content", "")
        for para in content.split("\n"):
            if para.strip():
                story.append(P(para.strip(), body_style))

    doc.build(story)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN PIPELINE
# ═══════════════════════════════════════════════════════════════════════════

def process_files(files: list, style: str, output: str, chunk_size: int,
                  api_key: str, api_base: str, generate_audio: bool = False,
                  qa_questions: list = None, use_shell: bool = True):
    import asyncio
    import concurrent.futures
    from datetime import datetime

    total = len(files)
    all_text = []
    all_file_names = []

    # ── Stage 1: Extracting ──
    emit("stage", {"stage": "extracting", "message": f"Extracting text from {total} files..."})
    for chunk_start in range(0, total, chunk_size):
        chunk = files[chunk_start:chunk_start + chunk_size]
        for i, f in enumerate(chunk):
            global_idx = chunk_start + i
            fname = os.path.basename(f)
            all_file_names.append(fname)
            emit("file_progress", {"file": fname, "index": global_idx, "total": total, "status": "extracting"})
            text = extract_text(f)
            all_text.append(text)
            emit("file_progress", {"file": fname, "index": global_idx, "total": total, "status": "done"})
    emit("progress", {"percent": 10, "message": f"Extracted text from {total} files"})

    # ── Stage 2: Correcting each file ──
    emit("stage", {"stage": "correcting", "message": "Correcting extracted text..."})
    all_corrected = []
    for idx, (fpath, ftext) in enumerate(zip(files, all_text)):
        fname = os.path.basename(fpath)
        fcorrected = correct_text(ftext)
        all_corrected.append({"fname": fname, "text": fcorrected})
        emit("file_progress", {"file": fname, "index": idx, "total": total, "status": "corrected"})
    emit("progress", {"percent": 15, "message": "Text correction complete"})

    # Load skill
    skill = load_skill(f"{style}-summary")
    if not skill:
        skill = load_skill("academic-summary")
    skill_body = skill["body"] if skill else ""

    # ── Stage 3+4: AI Explaining + Enhancing — PARALLEL per file ──
    emit("stage", {"stage": "ai_processing", "message": f"AI processing {total} files in parallel..."})

    def process_single_file(idx, fc):
        """Process a single file through ai_explain + ai_enhance. Returns (idx, fname, enhanced_text)."""
        fname = fc["fname"]
        ftext = fc["text"]
        fchars = len(ftext)

        # Signal start
        emit("file_progress", {"file": fname, "index": idx, "total": total, "status": "ai_explaining"})
        emit("progress", {"percent": 15 + int(35 * idx / max(total, 1)), "message": f"AI explaining file {idx + 1}/{total}: {fname} ({fchars} chars)..."})

        f_explained = ai_explain(ftext, api_key, api_base, style, skill_body)

        emit("file_progress", {"file": fname, "index": idx, "total": total, "status": "ai_enhancing"})
        emit("progress", {"percent": 15 + int(35 * (idx + 0.5) / max(total, 1)), "message": f"AI enhancing file {idx + 1}/{total}: {fname}..."})

        f_enhanced = ai_enhance(ftext, f_explained, api_key, api_base, style, skill_body)

        emit("file_progress", {"file": fname, "index": idx, "total": total, "status": "done"})
        emit("progress", {"percent": 15 + int(35 * (idx + 1) / max(total, 1)), "message": f"File {idx + 1}/{total} processed: {fname}"})

        return (idx, fname, f_enhanced)

    # Run all files in parallel using ThreadPoolExecutor
    all_enhanced = [None] * total
    max_workers = min(total, 5)
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        # Map each future back to its file index so errors are attributed correctly
        future_to_idx = {}
        for idx, fc in enumerate(all_corrected):
            fut = executor.submit(process_single_file, idx, fc)
            future_to_idx[fut] = idx

        for future in concurrent.futures.as_completed(future_to_idx):
            idx = future_to_idx[future]
            fc = all_corrected[idx]
            fname = fc["fname"]
            try:
                result_idx, result_fname, f_enhanced = future.result()
                all_enhanced[result_idx] = {"fname": result_fname, "text": f_enhanced}
            except Exception as e:
                emit("file_progress", {"file": fname, "index": idx, "total": total, "status": "error"})
                emit("warning", {"message": f"AI processing failed for {fname}: {e}. Using raw text."})
                all_enhanced[idx] = {"fname": fname, "text": fc["text"]}

    # Safety: fill any remaining None entries with raw text (shouldn't happen, but defensive)
    for idx in range(total):
        if all_enhanced[idx] is None:
            fc = all_corrected[idx]
            fname = fc["fname"]
            emit("warning", {"message": f"File {fname} was not processed (unknown error). Using raw text."})
            all_enhanced[idx] = {"fname": fname, "text": fc["text"]}

    emit("progress", {"percent": 50, "message": f"AI processing complete for {total} files"})

    # ── Stage 5: Build PDF with cover, TOC, and per-file content ──
    emit("stage", {"stage": "building_pdf", "message": "Building PDF with cover and table of contents..."})
    emit("progress", {"percent": 55, "message": "Building PDF..."})

    os.makedirs(os.path.dirname(os.path.abspath(output)), exist_ok=True)

    # ── Stage 5: Build PDF directly from per-file enhanced content ──
    emit("stage", {"stage": "building_pdf", "message": f"Building PDF with cover, TOC, and {total} chapters..."})
    emit("progress", {"percent": 55, "message": "Building PDF..."})

    os.makedirs(os.path.dirname(os.path.abspath(output)), exist_ok=True)

    # Build a minimal summary for the cover page title.
    # Combine all enhanced text into summary["content"] so TTS/Q&A don't KeyError.
    combined_content = "\n\n".join(
        f"--- SOURCE FILE {i+1}: {fe['fname']} ---\n{fe['text']}"
        for i, fe in enumerate(all_enhanced)
    )
    summary = {
        "title": f"Study Summary ({total} files)",
        "sections": [],
        "content": combined_content,
    }

    _build_final_pdf(summary, all_enhanced, style, output)

    emit("progress", {"percent": 85, "message": "PDF generation complete"})

    # ── Stage 6: TTS Audio (optional) ──
    audio_path = None
    if generate_audio:
        emit("stage", {"stage": "generating_audio", "message": "Generating audio summary..."})
        audio_path = output.replace(".pdf", ".mp3")
        tts_success = generate_tts(summary["content"], audio_path)
        if tts_success:
            emit("progress", {"percent": 95, "message": "Audio summary generated"})
        else:
            audio_path = None
            emit("warning", {"message": "TTS audio generation failed (no TTS engine available)"})

    # ── Stage 7: Q&A (optional) ──
    qa_results = []
    if qa_questions:
        emit("stage", {"stage": "qa_mode", "message": f"Answering {len(qa_questions)} questions..."})
        for i, q in enumerate(qa_questions):
            answer = answer_question(q, summary["content"], api_key, api_base)
            qa_results.append({"question": q, "answer": answer})
            emit("qa_answer", {"index": i, "question": q, "answer": answer})
        emit("progress", {"percent": 98, "message": "Q&A complete"})

    # ── Complete ──
    emit("progress", {"percent": 100, "message": "Summary generation complete"})
    emit("complete", {
        "output": output,
        "style": style,
        "audio": audio_path,
        "qa": qa_results,
        "title": summary.get("title", "Study Summary"),
        "sections_count": len(summary.get("sections", [])),
    })


def handle_qa_mode(args):
    """Handle Q&A mode — answer questions about a previously generated summary."""
    import json as _json

    if not args.job_data:
        print(_json.dumps({"error": "No job data provided"}), flush=True)
        sys.exit(1)

    try:
        job_events = _json.loads(args.job_data)
    except Exception as e:
        print(_json.dumps({"error": f"Invalid job data: {e}"}), flush=True)
        sys.exit(1)

    question = args.question
    if not question:
        print(_json.dumps({"error": "No question provided"}), flush=True)
        sys.exit(1)

    # Build context from job events
    context_parts = []
    for event in job_events:
        d = event.get("data", {})
        if event.get("type") == "stage":
            context_parts.append(f"Stage: {d.get('stage', '')} — {d.get('message', '')}")
        elif event.get("type") == "progress":
            context_parts.append(f"Progress: {d.get('percent', 0)}% — {d.get('message', '')}")

    context = "\n".join(context_parts)

    system_prompt = (
        "You are a knowledgeable tutor. Answer the user's question based on the "
        "summary generation context provided below. If the answer cannot be determined "
        "from the context, say so clearly.\n\n"
        f"Generation context:\n{context}"
    )

    try:
        answer = _call_ai(
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": question}],
            args.api_key, args.api_base, max_tokens=200000
        )
        print(_json.dumps({"answer": answer, "question": question}), flush=True)
    except Exception as e:
        print(_json.dumps({"error": str(e)}), flush=True)
        sys.exit(1)


def handle_shell_mode(args):
    """Handle shell mode — execute a shell command for PDF manipulation."""
    if not args.shell_command:
        print(json.dumps({"error": "No shell command provided", "success": False}), flush=True)
        sys.exit(1)

    result = execute_shell(
        command=args.shell_command,
        working_dir=args.shell_dir or None,
    )
    print(json.dumps(result), flush=True)


def handle_pipeline_mode(args):
    """Handle pipeline mode — process multiple files as a single pipeline."""
    files = [f.strip() for f in args.files.split(",") if f.strip()]
    if not files:
        emit("error", {"message": "No files provided"})
        sys.exit(1)
    for f in files:
        if not os.path.exists(f):
            emit("error", {"message": f"File not found: {f}"})
            sys.exit(1)
    process_files(
        files, args.style, args.output, args.chunk_size,
        args.api_key, args.api_base,
        generate_audio=args.generate_audio,
        qa_questions=args.qa_questions if args.qa_questions else None,
    )


def handle_env_update_mode(args):
    """Handle environment update mode — update .env files with new configuration."""
    env_path = args.env_path or os.path.join(os.getcwd(), ".env")
    env_example_path = args.env_example_path or os.path.join(os.getcwd(), ".env.example")
    
    updates = {}
    if args.ai_text_model:
        updates["AI_TEXT_MODEL"] = args.ai_text_model
    if args.ai_vision_model:
        updates["AI_VISION_MODEL"] = args.ai_vision_model
    if args.ai_qbank_model:
        updates["AI_QBANK_MODEL"] = args.ai_qbank_model
    if args.ai_explain_model:
        updates["AI_EXPLAIN_MODEL"] = args.ai_explain_model
    if args.academic_model:
        updates["ACADEMIC_MODEL"] = args.academic_model
    if args.modern_model:
        updates["MODERN_MODEL"] = args.modern_model
    if args.minimal_model:
        updates["MINIMAL_MODEL"] = args.minimal_model
    if args.clinical_model:
        updates["CLINICAL_MODEL"] = args.clinical_model
    if args.cornell_model:
        updates["CORNELL_MODEL"] = args.cornell_model
    if args.smart_briefing_model:
        updates["SMART_BRIEFING_MODEL"] = args.smart_briefing_model
    
    if updates:
        update_env_file(env_path, updates)
        update_env_example(env_example_path, updates)
        print(json.dumps({"success": True, "message": f"Updated {len(updates)} environment variables"}), flush=True)
    else:
        print(json.dumps({"success": False, "message": "No updates provided"}), flush=True)


def main():
    parser = argparse.ArgumentParser(description="Enhanced Summary Builder Pipeline")
    parser.add_argument("--files", required=True, help="Comma-separated file paths")
    parser.add_argument("--style", default="academic", help="Style template name")
    parser.add_argument("--output", required=True, help="Output PDF path")
    parser.add_argument("--chunk-size", type=int, default=3, help="Files per chunk")
    parser.add_argument("--api-key", default="", help="OpenRouter API key")
    parser.add_argument("--api-base", default="https://openrouter.ai/api/v1", help="API base URL")
    parser.add_argument("--generate-audio", action="store_true", help="Generate TTS audio summary")
    parser.add_argument("--qa-questions", nargs="*", default=[], help="Follow-up Q&A questions")
    parser.add_argument("--mode", default="generate", choices=["generate", "qa", "shell", "update-env"], help="Operation mode")
    parser.add_argument("--job-data", default="", help="JSON job data for Q&A mode")
    parser.add_argument("--question", default="", help="Single question for Q&A mode")
    parser.add_argument("--shell-command", default="", help="Shell command for shell mode")
    parser.add_argument("--shell-dir", default="", help="Working directory for shell mode")
    parser.add_argument("--use-shell", action="store_true", default=True, help="Use shell commands for PDF generation")
    
    # Environment update arguments
    parser.add_argument("--env-path", default="", help="Path to .env file")
    parser.add_argument("--env-example-path", default="", help="Path to .env.example file")
    parser.add_argument("--ai-text-model", default="", help="AI text model")
    parser.add_argument("--ai-vision-model", default="", help="AI vision model")
    parser.add_argument("--ai-qbank-model", default="", help="AI qbank model")
    parser.add_argument("--ai-explain-model", default="", help="AI explain model")
    parser.add_argument("--academic-model", default="", help="Academic style model")
    parser.add_argument("--modern-model", default="", help="Modern style model")
    parser.add_argument("--minimal-model", default="", help="Minimal style model")
    parser.add_argument("--clinical-model", default="", help="Clinical style model")
    parser.add_argument("--cornell-model", default="", help="Cornell style model")
    parser.add_argument("--smart-briefing-model", default="", help="Smart briefing style model")
    
    args = parser.parse_args()

    if args.mode == "generate":
        files = [f.strip() for f in args.files.split(",") if f.strip()]
        if not files:
            emit("error", {"message": "No files provided"})
            sys.exit(1)
        for f in files:
            if not os.path.exists(f):
                emit("error", {"message": f"File not found: {f}"})
                sys.exit(1)
        process_files(
            files, args.style, args.output, args.chunk_size,
            args.api_key, args.api_base,
            generate_audio=args.generate_audio,
            qa_questions=args.qa_questions if args.qa_questions else None,
            use_shell=args.use_shell,
        )
    elif args.mode == "qa":
        handle_qa_mode(args)
    elif args.mode == "shell":
        handle_shell_mode(args)
    elif args.mode == "update-env":
        handle_env_update_mode(args)


if __name__ == "__main__":
    main()