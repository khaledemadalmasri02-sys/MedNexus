#!/usr/bin/env python3
"""
Summary Builder — OpenMinis-inspired AI summary generation pipeline.

Features:
  - Skills system: modular SKILL.md-like prompt modules with metadata
  - Smart briefing mode: multi-source merge with actionable takeaways
  - TTS audio output: generate spoken summary
  - Q&A mode: follow-up questions about generated summary
  - Built-in Linux shell: execute shell commands for advanced PDF manipulation
"""

import argparse
import json
import os
import re
import sys
import csv
import io
import subprocess
import tempfile
from pathlib import Path


# ═══════════════════════════════════════════════════════════════════════════
# EVENT EMISSION
# ═══════════════════════════════════════════════════════════════════════════

def emit(event_type: str, data: dict):
    """Write a JSON-line event to stdout for the backend to stream."""
    payload = json.dumps({"type": event_type, **data})
    print(payload, flush=True)


# ═══════════════════════════════════════════════════════════════════════════
# SKILLS SYSTEM — Load SKILL.md modules from skills directory
# ═══════════════════════════════════════════════════════════════════════════

SKILLS_DIR = Path(__file__).parent / "skills"


def load_skill(skill_name: str) -> dict:
    """Load a skill module by name. Returns {name, description, body} or None."""
    skill_dir = SKILLS_DIR / skill_name
    skill_file = skill_dir / "SKILL.md"
    if not skill_file.exists():
        emit("warning", {"message": f"Skill not found: {skill_name}"})
        return None

    content = skill_file.read_text(encoding="utf-8")

    # Parse YAML frontmatter
    name = skill_name
    description = ""
    body = content

    if content.startswith("---"):
        parts = content.split("---", 2)
        if len(parts) >= 3:
            frontmatter = parts[1].strip()
            body = parts[2].strip()

            for line in frontmatter.split("\n"):
                line = line.strip()
                if line.startswith("name:"):
                    name = line.split(":", 1)[1].strip()
                elif line.startswith("description:"):
                    description = line.split(":", 1)[1].strip()
                elif description and not line.startswith("name:"):
                    # Multi-line description
                    description += " " + line.strip()

    return {"name": name, "description": description, "body": body}


def load_all_skills() -> dict:
    """Load all available skills from the skills directory."""
    skills = {}
    if not SKILLS_DIR.exists():
        return skills
    for skill_dir in sorted(SKILLS_DIR.iterdir()):
        if skill_dir.is_dir():
            skill = load_skill(skill_dir.name)
            if skill:
                skills[skill_dir.name] = skill
    return skills


def get_style_prompts(style: str) -> dict:
    """Get prompts for a style, using skills system if available."""
    skill = load_skill(f"{style}-summary")
    if skill:
        return {
            "explain": skill["body"],
            "enhance": skill["body"],
            "pdf_write": skill["body"],
            "system": skill["body"],
        }
    # Fallback for unknown styles
    generic = load_skill("academic-summary")
    if generic:
        return {
            "explain": generic["body"],
            "enhance": generic["body"],
            "pdf_write": generic["body"],
            "system": generic["body"],
        }
    return {"explain": "", "enhance": "", "pdf_write": "", "system": ""}


# ═══════════════════════════════════════════════════════════════════════════
# TEXT EXTRACTION (unchanged)
# ═══════════════════════════════════════════════════════════════════════════

def detect_type(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    mapping = {
        ".pdf": "pdf",
        ".png": "image", ".jpg": "image", ".jpeg": "image", ".webp": "image",
        ".pptx": "pptx", ".ppt": "pptx",
        ".xlsx": "spreadsheet", ".xls": "spreadsheet", ".csv": "spreadsheet",
        ".txt": "text", ".md": "text",
        ".docx": "docx",
    }
    return mapping.get(ext, "unknown")


def extract_text_pdf(file_path: str) -> str:
    import fitz
    doc = fitz.open(file_path)
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
    full_text = "\n".join(text_parts)
    doc.close()
    if len(full_text.strip()) < 100:
        return extract_pdf_ocr(file_path)
    return full_text


def extract_pdf_ocr(file_path: str) -> str:
    try:
        from pdf2image import convert_from_path
        import pytesseract
        images = convert_from_path(file_path, dpi=200)
        text_parts = []
        for img in images:
            text_parts.append(pytesseract.image_to_string(img))
        return "\n".join(text_parts)
    except Exception:
        return extract_text_pdf_fallback(file_path)


def extract_text_pdf_fallback(file_path: str) -> str:
    import fitz
    doc = fitz.open(file_path)
    text_parts = []
    for page in doc:
        blocks = page.get_text("blocks")
        for block in blocks:
            if len(block) >= 5 and block[4]:
                text_parts.append(block[4].strip())
    doc.close()
    return "\n".join(text_parts)


def extract_text_image(file_path: str) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        return pytesseract.image_to_string(img, config="--psm 6")
    except Exception:
        return ""


def extract_text_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_text.append(line)
        if slide_text:
            text_parts.append("\n".join(slide_text))
    return "\n\n".join(text_parts)


def extract_text_spreadsheet(file_path: str) -> str:
    import openpyxl
    ext = Path(file_path).suffix.lower()
    text_parts = []
    if ext == ".csv":
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                text_parts.append(" | ".join(row))
    else:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                text_parts.append(" | ".join(["" if c is None else str(c) for c in row]))
        wb.close()
    return "\n".join(text_parts)


def extract_text_plain(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def extract_text_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""


def extract_text(file_path: str) -> str:
    ftype = detect_type(file_path)
    handlers = {
        "pdf": extract_text_pdf,
        "image": extract_text_image,
        "pptx": extract_text_pptx,
        "spreadsheet": extract_text_spreadsheet,
        "text": extract_text_plain,
        "docx": extract_text_docx,
    }
    handler = handlers.get(ftype, extract_text_plain)
    try:
        return handler(file_path)
    except Exception as e:
        emit("file_error", {"file": file_path, "error": str(e)})
        return ""


def correct_text(text: str) -> str:
    text = re.sub(r"-\n", "", text)
    text = re.sub(r"([a-zA-Z,;:])\n([a-zA-Z])", r"\1 \2", text)
    text = re.sub(r"  +", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[^\x20-\x7E\n\r\t\u00A0-\uD7FF\uE000-\uFFFD]", "", text)
    lines = text.split("\n")
    merged = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            merged.append("")
            continue
        if merged and merged[-1]:
            prev = merged[-1]
            if (
                not prev.endswith((".", "!", "?", ":", ";", "-", "—"))
                and stripped
                and stripped[0].islower()
                and not prev.startswith(("•", "-", "*", "→", "#", "[", "Step", "Note"))
            ):
                merged[-1] = prev + " " + stripped
                continue
        merged.append(stripped)
    return "\n".join(merged)


# ═══════════════════════════════════════════════════════════════════════════
# AI API CALLS
# ═══════════════════════════════════════════════════════════════════════════

def _call_ai(messages: list, api_key: str, api_base: str, max_tokens: int = 4096) -> str:
    """Make an AI API call and return the response content."""
    import urllib.request

    payload = json.dumps({
        "model": "openrouter/owl-alpha",
        "messages": messages,
        "max_tokens": max_tokens,
    }).encode()

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}",
    }

    url = f"{api_base.rstrip('/')}/chat/completions"
    req = urllib.request.Request(url, data=payload, headers=headers)
    with urllib.request.urlopen(req, timeout=180) as resp:
        result = json.loads(resp.read())
        return result["choices"][0]["message"]["content"]


def ai_explain(text: str, api_key: str, api_base: str, style: str, skill_body: str) -> str:
    """Stage 1: AI explains key concepts using loaded skill prompts."""
    system_prompt = (
        f"{skill_body}\n\n"
        "Explain and elaborate on the following extracted study material. "
        "Add depth, clarify concepts, and make it study-ready."
    )
    user_prompt = f"{text[:6000]}"

    try:
        content = _call_ai(
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
            api_key, api_base, max_tokens=4096
        )
        return content
    except Exception as e:
        emit("warning", {"message": f"AI explanation failed: {e}. Using corrected text as-is."})
        return text


def ai_enhance(original_text: str, explained_text: str, api_key: str, api_base: str, style: str, skill_body: str) -> str:
    """Stage 2: AI enriches content using loaded skill prompts."""
    system_prompt = (
        f"{skill_body}\n\n"
        "Combine the original text and the AI explanation below into a single, "
        "enhanced study document. Add all enhancements, callouts, tables, and summaries "
        "as described in the skill instructions."
    )
    user_prompt = (
        "Original extracted text:\n\n"
        f"{original_text[:4000]}\n\n"
        "---\n\n"
        "AI-generated explanation:\n\n"
        f"{explained_text[:4000]}"
    )

    try:
        content = _call_ai(
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
            api_key, api_base, max_tokens=4096
        )
        return content
    except Exception as e:
        emit("warning", {"message": f"AI enhancement failed: {e}. Using explained text as-is."})
        return explained_text


def ai_write_pdf(enhanced_text: str, api_key: str, api_base: str, style: str, skill_body: str) -> dict:
    """Stage 3: AI writes complete PDF content using loaded skill prompts."""
    system_prompt = (
        f"{skill_body}\n\n"
        "Now write the COMPLETE PDF document content based on the enhanced content below. "
        "This will be directly rendered into a PDF. Include ALL sections as instructed in the skill. "
        "Be comprehensive and thorough — this is the entire document."
    )
    user_prompt = f"Enhanced study content:\n\n{enhanced_text[:6000]}"

    try:
        content = _call_ai(
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
            api_key, api_base, max_tokens=6000
        )
        return {
            "content": content,
            "sections": _parse_sections(content),
            "title": _extract_title(content),
        }
    except Exception as e:
        emit("warning", {"message": f"AI PDF writing failed: {e}. Falling back to template generation."})
        return {
            "content": enhanced_text,
            "sections": _parse_sections(enhanced_text),
            "title": _extract_title(enhanced_text),
        }


def _extract_title(text: str) -> str:
    for line in text.split("\n"):
        stripped = line.strip().lstrip("#").strip().lstrip("🎯📌💡⚡🔗📋🔬📊💊⚠️❓📚 ")
        if stripped and len(stripped) < 120:
            return stripped
    return "Study Summary"


def _parse_sections(text: str) -> list:
    sections = []
    pattern = re.compile(r"^#{1,3}\s+(.+)$", re.MULTILINE)
    matches = list(pattern.finditer(text))
    for i, m in enumerate(matches):
        name = m.group(1).strip()
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body = text[start:end].strip()
        sections.append({"name": name, "body": body})
    return sections


# ═══════════════════════════════════════════════════════════════════════════
# TTS AUDIO OUTPUT
# ═══════════════════════════════════════════════════════════════════════════

def generate_tts(text: str, output_path: str) -> bool:
    """Generate TTS audio from summary text. Returns True on success."""
    try:
        # Try edge-tts (free, built-in Windows/macOS/Linux TTS)
        import subprocess
        import sys

        # Create a condensed version for TTS (first 2000 chars of each section)
        tts_text = _prepare_tts_text(text)

        # Use edge-tts if available
        result = subprocess.run(
            [sys.executable, "-m", "edge_tts", "--voice", "en-US-AriaNeural",
             "--text", tts_text, "--write-media", output_path],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0:
            return True

    except (subprocess.TimeoutExpired, FileNotFoundError, Exception):
        pass

    # Fallback: try gTTS
    try:
        from gtts import gTTS
        tts_text = _prepare_tts_text(text)
        tts = gTTS(text=tts_text, lang="en")
        tts.save(output_path)
        return True
    except Exception:
        pass

    # Fallback: try espeak-ng
    try:
        result = subprocess.run(
            ["espeak-ng", "-w", output_path],
            input=text[:2000], capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return True
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception):
        pass

    return False


def _prepare_tts_text(text: str) -> str:
    """Prepare text for TTS — strip markdown, limit length, clean up."""
    # Remove markdown formatting
    clean = re.sub(r"[#*\-•→🎯💡📌🔗⚡📋🔬📊💊⚠️❓📚]\s*", "", text)
    clean = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", clean)  # Remove links
    clean = re.sub(r"`([^`]*)`", r"\1", clean)  # Remove code markers
    clean = re.sub(r"\n{2,}", ". ", clean)  # Replace double newlines with period
    clean = re.sub(r"\n", " ", clean)
    clean = re.sub(r"\s{2,}", " ", clean)
    # Limit to ~3000 chars for TTS
    return clean[:3000].strip()


# ═══════════════════════════════════════════════════════════════════════════
# Q&A MODE — Follow-up questions about generated summary
# ═══════════════════════════════════════════════════════════════════════════

def answer_question(question: str, summary_content: str, api_key: str, api_base: str) -> str:
    """Answer a follow-up question based on the generated summary content."""
    system_prompt = (
        "You are a knowledgeable tutor. Answer the user's question based ONLY on the "
        "provided summary content. If the answer is not in the summary, say so clearly. "
        "Be concise but thorough. Use examples from the summary when helpful.\n\n"
        "Summary content:\n\n"
        f"{summary_content[:6000]}"
    )

    try:
        answer = _call_ai(
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": question}],
            api_key, api_base, max_tokens=2048
        )
        return answer
    except Exception as e:
        return f"Error generating answer: {e}"


# ═══════════════════════════════════════════════════════════════════════════
# BUILT-IN LINUX SHELL — Execute shell commands for advanced PDF manipulation
# ═══════════════════════════════════════════════════════════════════════════

ALLOWED_SHELL_COMMANDS = {
    "pdftotext", "pdfinfo", "pdfunite", "pdftk", "ghostscript", "gs",
    "qpdf", "pdfseparate", "pdfjam", " enscript", "ps2pdf",
    "convert", "img2pdf", "exiftool", "pdffonts", "pdfimages",
    "mutool", "lp", "lpr", "cat", "grep", "sed", "awk", "head", "tail",
    "wc", "sort", "uniq", "tr", "cut", "paste", "join", "diff",
    "cp", "mv", "rm", "mkdir", "ls", "find", "file", "stat", "chmod",
    "tee", "xargs", "basename", "dirname", "mktemp", "date", "echo",
}

BLOCKED_SHELL_PATTERNS = [
    "rm -rf /", "rm -rf /*", "mkfs", "dd if=", ":(){:|:&};:",
    "chmod -R 777 /", "chown -R", "sudo", "su -", "> /dev/",
]


def execute_shell(command: str, working_dir: str = None, timeout: int = 30) -> dict:
    """
    Execute a shell command safely. Used for advanced PDF manipulation.
    Returns {stdout, stderr, exit_code, success}.
    """
    # Security check
    for blocked in BLOCKED_SHELL_PATTERNS:
        if blocked in command:
            return {"stdout": "", "stderr": f"Blocked: dangerous pattern '{blocked}'", "exit_code": 1, "success": False}

    # Check command is allowed
    cmd_parts = command.strip().split()
    if not cmd_parts:
        return {"stdout": "", "stderr": "Empty command", "exit_code": 1, "success": False}

    base_cmd = cmd_parts[0]
    if base_cmd not in ALLOWED_SHELL_COMMANDS:
        return {
            "stdout": "",
            "stderr": f"Command '{base_cmd}' not in allowed list. Allowed: {', '.join(sorted(ALLOWED_SHELL_COMMANDS))}",
            "exit_code": 1,
            "success": False,
        }

    try:
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=timeout,
            cwd=working_dir or os.getcwd(),
            env={**os.environ, "PATH": "/usr/local/bin:/usr/bin:/bin"},
        )
        return {
            "stdout": result.stdout.strip(),
            "stderr": result.stderr.strip(),
            "exit_code": result.returncode,
            "success": result.returncode == 0,
        }
    except subprocess.TimeoutExpired:
        return {"stdout": "", "stderr": f"Command timed out after {timeout}s", "exit_code": 1, "success": False}
    except Exception as e:
        return {"stdout": "", "stderr": str(e), "exit_code": 1, "success": False}


def shell_process_pdf(pdf_path: str, operation: str) -> dict:
    """
    Process a PDF using shell commands. Supported operations:
    - 'info': Get PDF metadata (page count, size, etc.)
    - 'extract_text': Extract all text from PDF
    - ' Compress': Compress PDF using ghostscript
    - 'merge': Merge with another PDF (requires second path)
    - 'split': Split PDF into individual pages
    - 'count_pages': Count total pages
    - 'get_fonts': List embedded fonts
    """
    operations = {
        "info": f"pdfinfo '{pdf_path}' 2>/dev/null || mutool info '{pdf_path}' 2>/dev/null",
        "extract_text": f"pdftotext '{pdf_path}' - 2>/dev/null || mutool draw -F text '{pdf_path}' 2>/dev/null",
        "compress": f"gs -sDEVICE=pdfwrite -dCompatibilityLevel=1.4 -dPDFSETTINGS=/ebook -dNOPAUSE -dBATCH -sOutputFile='{pdf_path}.compressed.pdf' '{pdf_path}' 2>/dev/null",
        "count_pages": f"pdfinfo '{pdf_path}' 2>/dev/null | grep Pages: | awk '{{print $2}}'",
        "get_fonts": f"pdffonts '{pdf_path}' 2>/dev/null",
        "split": f"mkdir -p '{pdf_path}.pages' && pdfseparate '{pdf_path}' '{pdf_path}.pages/page-%03d.pdf' 2>/dev/null",
    }

    cmd = operations.get(operation)
    if not cmd:
        return {"stdout": "", "stderr": f"Unknown operation: {operation}. Available: {', '.join(operations.keys())}", "exit_code": 1, "success": False}

    return execute_shell(command=cmd, timeout=60)


# ═══════════════════════════════════════════════════════════════════════════
# PDF GENERATION (unchanged ReportLab rendering)
# ═══════════════════════════════════════════════════════════════════════════

class AcademicStyle:
    name = "Academic"
    title_font = "Times-Bold"
    body_font = "Times-Roman"
    title_size = 22
    heading_size = 14
    body_size = 11
    title_color = (0.15, 0.15, 0.35)
    heading_color = (0.2, 0.2, 0.45)
    body_color = (0.1, 0.1, 0.1)
    accent_color = (0.25, 0.35, 0.55)
    line_color = (0.7, 0.7, 0.75)
    bg_tint = (0.97, 0.97, 0.99)
    heading_bg = (0.92, 0.93, 0.97)
    margin = 72
    use_heading_bg = True
    use_section_lines = True
    title_underline = True


class ModernStyle:
    name = "Modern"
    title_font = "Helvetica-Bold"
    body_font = "Helvetica"
    title_size = 26
    heading_size = 16
    body_size = 11
    title_color = (0.1, 0.1, 0.15)
    heading_color = (0.0, 0.55, 0.75)
    body_color = (0.15, 0.15, 0.15)
    accent_color = (0.0, 0.55, 0.75)
    line_color = (0.0, 0.55, 0.75)
    bg_tint = (0.98, 0.99, 1.0)
    heading_bg = (0.9, 0.96, 1.0)
    margin = 55
    use_heading_bg = True
    use_section_lines = False
    title_underline = False
    title_bar = True


class MinimalStyle:
    name = "Minimal"
    title_font = "Helvetica"
    body_font = "Helvetica"
    title_size = 20
    heading_size = 13
    body_size = 10
    title_color = (0.2, 0.2, 0.2)
    heading_color = (0.35, 0.35, 0.35)
    body_color = (0.25, 0.25, 0.25)
    accent_color = (0.5, 0.5, 0.5)
    line_color = (0.85, 0.85, 0.85)
    bg_tint = (1.0, 1.0, 1.0)
    heading_bg = None
    margin = 85
    use_heading_bg = False
    use_section_lines = False
    title_underline = True
    minimal_dots = True


class ClinicalStyle:
    name = "Clinical"
    title_font = "Helvetica-Bold"
    body_font = "Helvetica"
    title_size = 28
    heading_size = 12
    body_size = 10
    title_color = (0.102, 0.322, 0.463)
    heading_color = (0.102, 0.322, 0.463)
    body_color = (0.129, 0.129, 0.129)
    accent_color = (0.102, 0.322, 0.463)
    line_color = (0.741, 0.741, 0.741)
    bg_tint = (1.0, 1.0, 1.0)
    heading_bg = None
    margin = 72
    use_heading_bg = False
    use_section_lines = False
    title_underline = False
    clinical_boxes = True


class CornellStyle:
    name = "Cornell"
    title_font = "Helvetica-Bold"
    body_font = "Helvetica"
    title_size = 22
    heading_size = 13
    body_size = 10
    title_color = (0.15, 0.15, 0.25)
    heading_color = (0.2, 0.3, 0.5)
    body_color = (0.1, 0.1, 0.1)
    accent_color = (0.2, 0.3, 0.5)
    line_color = (0.7, 0.7, 0.75)
    bg_tint = (0.99, 0.99, 1.0)
    heading_bg = (0.94, 0.95, 0.98)
    margin = 72
    use_heading_bg = True
    use_section_lines = True
    title_underline = True
    cornell_layout = True


class SmartBriefingStyle:
    name = "Smart Briefing"
    title_font = "Helvetica-Bold"
    body_font = "Helvetica"
    title_size = 24
    heading_size = 15
    body_size = 11
    title_color = (0.05, 0.15, 0.3)
    heading_color = (0.0, 0.45, 0.65)
    body_color = (0.15, 0.15, 0.15)
    accent_color = (0.0, 0.55, 0.45)
    line_color = (0.7, 0.75, 0.8)
    bg_tint = (0.98, 0.99, 1.0)
    heading_bg = (0.92, 0.97, 1.0)
    margin = 65
    use_heading_bg = True
    use_section_lines = False
    title_underline = True


STYLES = {
    "academic": AcademicStyle,
    "modern": ModernStyle,
    "minimal": MinimalStyle,
    "clinical": ClinicalStyle,
    "cornell": CornellStyle,
    "smart-briefing": SmartBriefingStyle,
}


def generate_pdf(summary: dict, style_name: str, output_path: str):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.colors import Color
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    HRFlowable, Frame, PageTemplate, Table,
                                    TableStyle)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    style = STYLES.get(style_name, AcademicStyle)
    page_w, page_h = A4

    def draw_background(canvas, doc):
        canvas.saveState()
        if style.bg_tint:
            canvas.setFillColor(Color(*style.bg_tint))
            canvas.rect(0, 0, page_w, page_h, fill=1, stroke=0)
        if getattr(style, 'title_bar', False) and doc.page == 1:
            canvas.setFillColor(Color(*style.accent_color))
            canvas.rect(0, page_h - 8*mm, page_w, 8*mm, fill=1, stroke=0)
        canvas.restoreState()

    doc = SimpleDocTemplate(
        output_path, pagesize=A4,
        leftMargin=style.margin, rightMargin=style.margin,
        topMargin=style.margin, bottomMargin=style.margin,
    )
    frame = Frame(style.margin, style.margin,
                  page_w - 2 * style.margin, page_h - 2 * style.margin, id='normal')
    template = PageTemplate(id='main', frames=frame, onPage=draw_background)
    doc.addPageTemplates([template])

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "SummaryTitle", parent=styles["Title"],
        fontName=style.title_font, fontSize=style.title_size,
        textColor=style.title_color, spaceAfter=6, leading=style.title_size * 1.3,
    )
    heading_style = ParagraphStyle(
        "SummaryHeading", parent=styles["Heading2"],
        fontName=style.title_font, fontSize=style.heading_size,
        textColor=style.heading_color, spaceBefore=14, spaceAfter=4,
        leading=style.heading_size * 1.3,
    )
    body_style = ParagraphStyle(
        "SummaryBody", parent=styles["Normal"],
        fontName=style.body_font, fontSize=style.body_size,
        textColor=style.body_color, spaceAfter=6, leading=style.body_size * 1.5,
    )

    story = []
    title = summary.get("title", "Study Summary")

    if getattr(style, 'title_bar', False):
        bar_title_style = ParagraphStyle("BarTitle", parent=title_style, textColor=(1, 1, 1), spaceAfter=2)
        story.append(Spacer(1, 2*mm))
        story.append(Paragraph(title, bar_title_style))
        story.append(Spacer(1, 4*mm))
    else:
        story.append(Paragraph(title, title_style))
        if getattr(style, 'title_underline', False):
            story.append(HRFlowable(
                width="100%", thickness=2 if style_name == "modern" else 1,
                color=style.accent_color if getattr(style, 'title_bar', False) else style.line_color,
                spaceAfter=12, spaceBefore=4,
            ))
        else:
            story.append(Spacer(1, 6*mm))

    sections = summary.get("sections", [])

    if getattr(style, 'cornell_layout', False) and sections:
        story.append(Paragraph("Cue Column / Questions", heading_style))
        story.append(HRFlowable(width="100%", thickness=1, color=style.line_color, spaceAfter=8))
        for section in sections:
            story.append(Paragraph(f"• {section['name']}", body_style))
        story.append(Spacer(1, 8*mm))
        story.append(HRFlowable(width="100%", thickness=2, color=style.accent_color, spaceAfter=8))
        story.append(Paragraph("Notes", heading_style))
        for section in sections:
            story.append(Paragraph(section["name"], heading_style))
            for para in section["body"].split("\n"):
                para = para.strip()
                if para:
                    story.append(Paragraph(para, body_style))
            story.append(Spacer(1, 4*mm))
        story.append(Spacer(1, 8*mm))
        story.append(HRFlowable(width="100%", thickness=2, color=style.accent_color, spaceAfter=8))
        summary_style = ParagraphStyle("CornellSummary", parent=body_style,
                                       fontSize=style.body_size - 1, textColor=style.heading_color)
        story.append(Paragraph("Summary", heading_style))
        full_content = summary.get("content", "")
        summary_text = full_content[:500] + ("..." if len(full_content) > 500 else "")
        story.append(Paragraph(summary_text, summary_style))
    elif sections:
        for idx, section in enumerate(sections):
            if getattr(style, 'use_heading_bg', False) and style.heading_bg:
                heading_table = Table(
                    [[Paragraph(section["name"], ParagraphStyle(
                        "BgHeading", parent=heading_style, textColor=style.heading_color,
                        spaceBefore=0, spaceAfter=0, leading=style.heading_size * 1.3,
                    ))]],
                    colWidths=[doc.width],
                )
                heading_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, -1), Color(*style.heading_bg)),
                    ('LEFTPADDING', (0, 0), (-1, -1), 8),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                    ('TOPPADDING', (0, 0), (-1, -1), 4),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                    ('ROUNDEDCORNERS', [4, 4, 4, 4]),
                ]))
                story.append(heading_table)
                story.append(Spacer(1, 4*mm))
            elif getattr(style, 'clinical_boxes', False):
                clinical_heading = Table(
                    [[Paragraph(section["name"], ParagraphStyle(
                        "ClinHeading", parent=heading_style, textColor=(1, 1, 1),
                        spaceBefore=0, spaceAfter=0, leading=style.heading_size * 1.3,
                    ))]],
                    colWidths=[doc.width],
                )
                clinical_heading.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, -1), Color(*style.accent_color)),
                    ('LEFTPADDING', (0, 0), (-1, -1), 10),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 10),
                    ('TOPPADDING', (0, 0), (-1, -1), 5),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
                ]))
                story.append(clinical_heading)
                story.append(Spacer(1, 4*mm))
            else:
                story.append(Paragraph(section["name"], heading_style))

            if getattr(style, 'use_section_lines', False):
                story.append(HRFlowable(
                    width="100%", thickness=0.5, color=style.line_color,
                    spaceAfter=4, spaceBefore=2,
                ))

            for para in section["body"].split("\n"):
                para = para.strip()
                if para:
                    if getattr(style, 'minimal_dots', False):
                        story.append(Paragraph(f"  • {para}", body_style))
                    else:
                        story.append(Paragraph(para, body_style))

            if getattr(style, 'clinical_boxes', False):
                box_content = []
                for para in section["body"].split("\n"):
                    if para.strip():
                        box_content.append(Paragraph(para.strip(), body_style))
                if box_content:
                    inner_table = Table([[e] for e in box_content], colWidths=[doc.width - 10*mm])
                    inner_table.setStyle(TableStyle([
                        ('LEFTPADDING', (0, 0), (-1, -1), 8),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                        ('TOPPADDING', (0, 0), (-1, -1), 2),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
                        ('BOX', (0, 0), (-1, -1), 0.5, Color(*style.line_color)),
                    ]))
                    story.append(inner_table)

            story.append(Spacer(1, 4 * mm))
    else:
        content = summary.get("content", "")
        for para in content.split("\n"):
            if para.strip():
                story.append(Paragraph(para.strip(), body_style))

    doc.build(story)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN PIPELINE
# ═══════════════════════════════════════════════════════════════════════════

def process_files(files: list, style: str, output: str, chunk_size: int,
                  api_key: str, api_base: str, generate_audio: bool = False,
                  qa_questions: list = None):
    total = len(files)
    all_text = []
    all_file_names = []

    # ── Stage 1: Extracting ──
    for chunk_start in range(0, total, chunk_size):
        chunk = files[chunk_start:chunk_start + chunk_size]
        chunk_num = (chunk_start // chunk_size) + 1
        total_chunks = (total + chunk_size - 1) // chunk_size
        emit("stage", {
            "stage": "extracting",
            "message": f"Processing files {chunk_start + 1}-{min(chunk_start + chunk_size, total)} of {total}",
            "chunk": chunk_num,
            "total_chunks": total_chunks,
        })

        for i, f in enumerate(chunk):
            global_idx = chunk_start + i
            fname = os.path.basename(f)
            all_file_names.append(fname)
            emit("file_progress", {"file": fname, "index": global_idx, "total": total, "status": "extracting"})
            text = extract_text(f)
            all_text.append(text)
            emit("file_progress", {"file": fname, "index": global_idx, "total": total, "status": "done"})

    # ── Stage 2: Correcting ──
    emit("stage", {"stage": "correcting", "message": "Correcting extracted text..."})
    combined = "\n\n".join(all_text)
    corrected = correct_text(combined)
    emit("progress", {"percent": 15, "message": "Text correction complete"})

    # Load skill for this style
    skill = load_skill(f"{style}-summary")
    if not skill:
        skill = load_skill("academic-summary")
    skill_body = skill["body"] if skill else ""

    # For smart-briefing, include file names in the context
    if style == "smart-briefing" and all_file_names:
        file_context = "\n\n---\n\nSource files:\n" + "\n".join(f"- {n}" for n in all_file_names)
        corrected = corrected + file_context

    # ── Stage 3: AI Explaining ──
    emit("stage", {"stage": "ai_explaining", "message": "AI explaining key concepts..."})
    explained = ai_explain(corrected, api_key, api_base, style, skill_body)
    emit("progress", {"percent": 30, "message": "AI explanation complete"})

    # ── Stage 4: AI Enhancing ──
    emit("stage", {"stage": "ai_enhancing", "message": "AI enhancing content..."})
    enhanced = ai_enhance(corrected, explained, api_key, api_base, style, skill_body)
    emit("progress", {"percent": 50, "message": "AI enhancement complete"})

    # ── Stage 5: AI Writing PDF Content ──
    emit("stage", {"stage": "ai_writing_pdf", "message": "AI writing PDF content..."})
    summary = ai_write_pdf(enhanced, api_key, api_base, style, skill_body)
    emit("progress", {"percent": 70, "message": "AI PDF content written"})

    # ── Stage 6: Building PDF ──
    emit("stage", {"stage": "building_pdf", "message": f"Building {style} PDF..."})
    os.makedirs(os.path.dirname(os.path.abspath(output)), exist_ok=True)
    generate_pdf(summary, style, output)

    # Use shell to get PDF info
    shell_result = shell_process_pdf(output, "info")
    if shell_result["success"]:
        emit("info", {"message": f"PDF created: {shell_result['stdout'][:200]}"})

    emit("progress", {"percent": 85, "message": "PDF generation complete"})

    # ── Stage 7: TTS Audio (optional) ──
    audio_path = None
    if generate_audio:
        emit("stage", {"stage": "generating_audio", "message": "Generating audio summary..."})
        audio_path = output.replace(".pdf", ".mp3")
        tts_success = generate_tts(summary["content"], audio_path)
        if tts_success:
            emit("progress", {"percent": 95, "message": "Audio summary generated"})
        else:
            audio_path = None
            emit("warning", {"message": "TTS audio generation failed (no TTS engine available)"})

    # ── Stage 8: Q&A (optional) ──
    qa_results = []
    if qa_questions:
        emit("stage", {"stage": "qa_mode", "message": f"Answering {len(qa_questions)} questions..."})
        for i, q in enumerate(qa_questions):
            answer = answer_question(q, summary["content"], api_key, api_base)
            qa_results.append({"question": q, "answer": answer})
            emit("qa_answer", {"index": i, "question": q, "answer": answer})
        emit("progress", {"percent": 98, "message": "Q&A complete"})

    # ── Complete ──
    emit("progress", {"percent": 100, "message": "Summary generation complete"})
    emit("complete", {
        "output": output,
        "style": style,
        "audio": audio_path,
        "qa": qa_results,
        "title": summary.get("title", "Study Summary"),
        "sections_count": len(summary.get("sections", [])),
    })


def handle_qa_mode(args):
    """Handle Q&A mode — answer questions about a previously generated summary."""
    import json as _json

    if not args.job_data:
        print(_json.dumps({"error": "No job data provided"}), flush=True)
        sys.exit(1)

    try:
        job_events = _json.loads(args.job_data)
    except Exception as e:
        print(_json.dumps({"error": f"Invalid job data: {e}"}), flush=True)
        sys.exit(1)

    # Find the complete event to get the summary content
    summary_content = ""
    for event in job_events:
        if event.get("type") == "complete":
            # The complete event has the output path; we need the actual content
            # Look for it in the events
            pass
        if event.get("type") == "stage" and event.get("data", {}).get("stage") == "ai_writing_pdf":
            # The content was written in a previous stage; reconstruct from events
            pass

    # For Q&A, we need the summary content. It's stored in the job events.
    # We'll extract it from the last known content-bearing event.
    # In a production system, this would be stored in a database.
    # For now, we pass the question and let the AI answer based on available context.

    question = args.question
    if not question:
        print(_json.dumps({"error": "No question provided"}), flush=True)
        sys.exit(1)

    # Build context from job events
    context_parts = []
    for event in job_events:
        d = event.get("data", {})
        if event.get("type") == "stage":
            context_parts.append(f"Stage: {d.get('stage', '')} — {d.get('message', '')}")
        elif event.get("type") == "progress":
            context_parts.append(f"Progress: {d.get('percent', 0)}% — {d.get('message', '')}")

    context = "\n".join(context_parts)

    system_prompt = (
        "You are a knowledgeable tutor. Answer the user's question based on the "
        "summary generation context provided below. If the answer cannot be determined "
        "from the context, say so clearly.\n\n"
        f"Generation context:\n{context}"
    )

    try:
        answer = _call_ai(
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": question}],
            args.api_key, args.api_base, max_tokens=2048
        )
        print(_json.dumps({"answer": answer, "question": question}), flush=True)
    except Exception as e:
        print(_json.dumps({"error": str(e)}), flush=True)
        sys.exit(1)


def handle_shell_mode(args):
    """Handle shell mode — execute a shell command for PDF manipulation."""
    if not args.shell_command:
        print(json.dumps({"error": "No shell command provided", "success": False}), flush=True)
        sys.exit(1)

    result = execute_shell(
        command=args.shell_command,
        working_dir=args.shell_dir or None,
    )
    print(json.dumps(result), flush=True)


def handle_pipeline_mode(args):
    """Handle pipeline mode — process multiple files as a single pipeline."""
    files = [f.strip() for f in args.files.split(",") if f.strip()]
    if not files:
        emit("error", {"message": "No files provided"})
        sys.exit(1)
    for f in files:
        if not os.path.exists(f):
            emit("error", {"message": f"File not found: {f}"})
            sys.exit(1)
    process_files(
        files, args.style, args.output, args.chunk_size,
        args.api_key, args.api_base,
        generate_audio=args.generate_audio,
        qa_questions=args.qa_questions if args.qa_questions else None,
    )


def main():
    parser = argparse.ArgumentParser(description="Summary Builder Pipeline")
    parser.add_argument("--files", required=True, help="Comma-separated file paths")
    parser.add_argument("--style", default="academic", help="Style template name")
    parser.add_argument("--output", required=True, help="Output PDF path")
    parser.add_argument("--chunk-size", type=int, default=3, help="Files per chunk")
    parser.add_argument("--api-key", default="", help="OpenRouter API key")
    parser.add_argument("--api-base", default="https://openrouter.ai/api/v1", help="API base URL")
    parser.add_argument("--generate-audio", action="store_true", help="Generate TTS audio summary")
    parser.add_argument("--qa-questions", nargs="*", default=[], help="Follow-up Q&A questions")
    parser.add_argument("--mode", default="generate", choices=["generate", "qa", "shell"], help="Operation mode")
    parser.add_argument("--job-data", default="", help="JSON job data for Q&A mode")
    parser.add_argument("--question", default="", help="Single question for Q&A mode")
    parser.add_argument("--shell-command", default="", help="Shell command for shell mode")
    parser.add_argument("--shell-dir", default="", help="Working directory for shell mode")
    args = parser.parse_args()

    if args.mode == "generate":
        files = [f.strip() for f in args.files.split(",") if f.strip()]
        if not files:
            emit("error", {"message": "No files provided"})
            sys.exit(1)
        for f in files:
            if not os.path.exists(f):
                emit("error", {"message": f"File not found: {f}"})
                sys.exit(1)
        process_files(
            files, args.style, args.output, args.chunk_size,
            args.api_key, args.api_base,
            generate_audio=args.generate_audio,
            qa_questions=args.qa_questions if args.qa_questions else None,
        )
    elif args.mode == "qa":
        handle_qa_mode(args)
    elif args.mode == "shell":
        handle_shell_mode(args)
    elif args.mode == "pipeline":
        handle_pipeline_mode(args)
    else:
        emit("error", {"message": f"Unknown mode: {args.mode}"})
        sys.exit(1)


if __name__ == "__main__":
    main()
